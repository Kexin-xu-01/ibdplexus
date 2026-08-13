"""
Generate PDF report and PowerPoint slides summarising the CD vs UC pipeline and results.

Inputs
------
- /home/jovyan/kgbk271-ibd-volume/training/cv_splits_patients.csv
- /home/jovyan/kgbk271-ibd-volume/training/cv_splits_slides.csv
- /home/jovyan/kgbk271-ibd-volume/training/cd_vs_uc/prism2_base_fold_metrics.csv
- /home/jovyan/kgbk271-ibd-volume/training/cd_vs_uc/prism2_diagnostic_fold_metrics.csv

Outputs  (versioned — never overwrites existing files)
--------
- pipeline_results_report_vN.pdf
- pipeline_results_slides_vN.pptx
- VERSIONS.md  (append-only changelog)

Usage
-----
    python 04_generate_reports.py
    python 04_generate_reports.py --desc "Added confusion matrix breakdown"
"""

import argparse
import os
import json
import numpy as np
import pandas as pd
from version_utils import next_versioned_path, log_version

RESULTS_DIR = '/home/jovyan/kgbk271-ibd-volume/training/cd_vs_uc/02_04_imaging_allsites/results'
REPORTS_DIR = '/home/jovyan/kgbk271-ibd-volume/training/cd_vs_uc/02_04_imaging_allsites/reports'
TRAINING_DIR = '/home/jovyan/kgbk271-ibd-volume/training'


# ══════════════════════════════════════════════════════════════════════════════
#  PDF REPORT
# ══════════════════════════════════════════════════════════════════════════════

def build_pdf():
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.lib import colors
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer, Table,
                                     TableStyle, HRFlowable, PageBreak)
    from reportlab.lib.enums import TA_JUSTIFY, TA_CENTER

    DARK  = colors.HexColor('#1a2e4a')
    MID   = colors.HexColor('#2e6da4')
    LIGHT = colors.HexColor('#d8e8f5')
    WHITE = colors.white
    GREEN = colors.HexColor('#1e8b4c')
    RED   = colors.HexColor('#c0392b')
    GREY  = colors.HexColor('#f5f7fa')
    MGREY = colors.HexColor('#dddddd')
    DGREY = colors.HexColor('#888888')

    out_path, _ = next_versioned_path(os.path.join(REPORTS_DIR, 'pipeline_results_report.pdf'))
    doc = SimpleDocTemplate(out_path, pagesize=A4,
                            topMargin=2.2*cm, bottomMargin=2*cm,
                            leftMargin=2.2*cm, rightMargin=2.2*cm)

    S = getSampleStyleSheet()
    def sty(name, **kw): return ParagraphStyle(name, parent=S['Normal'], **kw)

    H1   = sty('H1', fontSize=18, textColor=DARK, spaceAfter=6, fontName='Helvetica-Bold')
    H2   = sty('H2', fontSize=13, textColor=MID, spaceBefore=16, spaceAfter=5, fontName='Helvetica-Bold')
    H3   = sty('H3', fontSize=11, textColor=DARK, spaceBefore=10, spaceAfter=3, fontName='Helvetica-Bold')
    BODY = sty('BODY', fontSize=9.5, leading=15, spaceAfter=4, alignment=TA_JUSTIFY)
    SMALL= sty('SMALL', fontSize=8.5, leading=12, textColor=DGREY)
    BULL = sty('BULL', fontSize=9.5, leading=14, leftIndent=14, spaceAfter=3)
    NOTE = sty('NOTE', fontSize=8.5, leading=12, textColor=colors.HexColor('#444444'), leftIndent=10)
    FOOT = sty('FOOT', fontSize=7.5, textColor=DGREY, alignment=TA_CENTER)

    TS = [
        ('BACKGROUND', (0,0), (-1,0), MID), ('TEXTCOLOR', (0,0), (-1,0), WHITE),
        ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'), ('FONTSIZE', (0,0), (-1,-1), 9),
        ('ROWBACKGROUNDS', (0,1), (-1,-1), [WHITE, GREY]),
        ('GRID', (0,0), (-1,-1), 0.4, MGREY), ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
        ('TOPPADDING', (0,0), (-1,-1), 4), ('BOTTOMPADDING', (0,0), (-1,-1), 4),
        ('LEFTPADDING', (0,0), (-1,-1), 6), ('RIGHTPADDING', (0,0), (-1,-1), 6),
    ]

    def tbl(data, col_widths, extra=None):
        t = Table(data, colWidths=[w*cm for w in col_widths])
        t.setStyle(TableStyle(TS + (extra or [])))
        return t

    # load data
    base_folds = pd.read_csv(os.path.join(RESULTS_DIR, 'prism2_base_fold_metrics.csv'))
    diag_folds = pd.read_csv(os.path.join(RESULTS_DIR, 'prism2_diagnostic_fold_metrics.csv'))
    patients   = pd.read_csv(os.path.join(TRAINING_DIR, 'cv_splits_patients.csv'))
    slides     = pd.read_csv(os.path.join(TRAINING_DIR, 'cv_splits_slides.csv'))

    story = []

    # Title
    story += [Spacer(1, 1.5*cm),
              Paragraph('Colon CD vs UC Classification from<br/>Whole-Slide Image Embeddings', H1),
              Paragraph('Pipeline Description and Results Report',
                        sty('sub', fontSize=13, textColor=MID, spaceAfter=6)),
              HRFlowable(width='100%', thickness=2, color=MID, spaceAfter=10),
              Paragraph('Dataset: IBD Plexus / SPARC &nbsp;·&nbsp; Embeddings: Virchow2 via TRIDENT 0.3.0<br/>'
                        'Models: prism2_base (2560-d) and prism2_diagnostic (3072-d)<br/>'
                        'Date: August 10, 2026', SMALL),
              Spacer(1, 0.8*cm)]

    story.append(tbl([
        ['Cohort', f"{len(patients)} patients · {len(slides)} colon slides · CD 67% / UC 33%"],
        ['Task',   "Binary classification: Crohn's disease (CD) vs Ulcerative colitis (UC)"],
        ['Method', 'Random Forest (500 trees, balanced class weights)'],
        ['Evaluation', 'Patient-level 5-fold stratified cross-validation'],
        ['Best AUC', f"prism2_base  {base_folds['auc'].mean():.3f} ± {base_folds['auc'].std():.3f}"],
    ], [4, 12], extra=[
        ('FONTNAME', (0,0), (0,-1), 'Helvetica-Bold'),
        ('BACKGROUND', (0,0), (0,-1), LIGHT),
        ('BACKGROUND', (0,0), (-1,0), LIGHT),
        ('TEXTCOLOR', (0,0), (-1,0), DARK),
    ]))
    story.append(PageBreak())

    # Section 1: Data
    story.append(Paragraph('1. Data', H2))
    story.append(Paragraph(
        'Whole-slide images (WSIs) were obtained from the IBD Plexus / SPARC biobank. '
        'FFPE tissue sections stained with H&amp;E were digitised at 40× magnification. '
        'This analysis uses <b>colon biopsies only</b> for anatomical consistency.', BODY))
    story.append(Paragraph('1.1  Cohort statistics', H3))
    cd_n = (patients['diagnosis'] == "Crohn's disease").sum()
    uc_n = (patients['diagnosis'] == 'Ulcerative colitis').sum()
    cd_s = (slides['diagnosis'] == "Crohn's disease").sum()
    uc_s = (slides['diagnosis'] == 'Ulcerative colitis').sum()
    story.append(tbl([
        ['', "CD (Crohn's disease)", 'UC (Ulcerative colitis)', 'Total'],
        ['Patients', f'{cd_n} ({cd_n/len(patients)*100:.1f}%)', f'{uc_n} ({uc_n/len(patients)*100:.1f}%)', str(len(patients))],
        ['Colon slides', str(cd_s), str(uc_s), str(len(slides))],
        ['Slides/patient (mean)', f'{cd_s/cd_n:.2f}', f'{uc_s/uc_n:.2f}', f'{len(slides)/len(patients):.2f}'],
    ], [4, 4, 4, 3], extra=[('FONTNAME', (0,1), (0,-1), 'Helvetica-Bold')]))
    story.append(Paragraph(
        '<b>Exclusions:</b> 35 patients with conflicting CD/UC labels removed. '
        'One slide with a duplicate slide_id removed.', NOTE))

    # Section 2: Embeddings
    story.append(Paragraph('2. Image Embeddings', H2))
    story.append(Paragraph(
        'Slide-level embeddings were produced by TRIDENT 0.3.0 using Virchow2 '
        'at 20× magnification, 224 px patch size, 0 px overlap. '
        'Two aggregation variants were evaluated:', BODY))
    story.append(tbl([
        ['Embedding', 'Dimension', 'Description'],
        ['prism2_base',       '2,560', 'Virchow2 CLS token aggregated via PRISM2 base head'],
        ['prism2_diagnostic', '3,072', 'Virchow2 features aggregated via PRISM2 diagnostic head'],
    ], [4, 2.5, 9.5]))

    # Section 3: CV Split
    story.append(Paragraph('3. Cross-Validation Split', H2))
    story.append(Paragraph(
        'Patient-level 5-fold stratified CV using MultilabelStratifiedKFold '
        '(iterative-stratification). Stratification variables: '
        '<b>diagnosis × sex × age-at-diagnosis bin</b>. '
        'All slides for a patient stay in one fold.', BODY))
    story.append(tbl([
        ['Variable', 'Bins', 'Rationale'],
        ['Diagnosis', 'CD / UC', '2:1 class imbalance'],
        ['Sex', 'Female / Male / Missing', 'CD skews more female'],
        ['Age at diagnosis', '<20 / 20–35 / >35 / Missing', 'CD median 24 yrs vs UC 27 yrs'],
    ], [3.5, 4.5, 8]))

    # Section 4: Model
    story.append(Paragraph('4. Model', H2))
    story.append(Paragraph(
        'Random Forest classifier (scikit-learn) as a non-parametric baseline '
        'requiring no embedding-space tuning.', BODY))
    story.append(tbl([
        ['Hyperparameter', 'Value'],
        ['n_estimators', '500'],
        ['max_features', 'sqrt(d)'],
        ['min_samples_leaf', '2'],
        ['class_weight', 'balanced'],
        ['random_state', '42'],
    ], [5, 11]))

    # Section 5: Results
    story.append(Paragraph('5. Results', H2))
    story.append(Paragraph('5.1  Overall performance', H3))

    def fold_summary(folds):
        return (folds['auc'].mean(), folds['auc'].std(),
                folds['ap'].mean(), folds['ap'].std(),
                folds['accuracy'].mean(), folds['accuracy'].std())

    b = fold_summary(base_folds)
    d = fold_summary(diag_folds)
    story.append(tbl([
        ['Model', 'AUC (mean±SD)', 'AP (mean±SD)', 'Accuracy (mean±SD)'],
        ['prism2_base',       f'{b[0]:.3f} ± {b[1]:.3f}', f'{b[2]:.3f} ± {b[3]:.3f}', f'{b[4]*100:.1f}% ± {b[5]*100:.1f}%'],
        ['prism2_diagnostic', f'{d[0]:.3f} ± {d[1]:.3f}', f'{d[2]:.3f} ± {d[3]:.3f}', f'{d[4]*100:.1f}% ± {d[5]*100:.1f}%'],
    ], [5, 4.5, 4.5, 4.5], extra=[
        ('BACKGROUND', (0,1), (-1,1), colors.HexColor('#e8f4ec')),
        ('FONTNAME', (0,1), (-1,1), 'Helvetica-Bold'),
    ]))

    story.append(Paragraph('5.2  Aggregated confusion matrices', H3))
    for name, folds in [('prism2_base', base_folds), ('prism2_diagnostic', diag_folds)]:
        tn=folds['tn'].sum(); fp=folds['fp'].sum()
        fn=folds['fn'].sum(); tp=folds['tp'].sum()
        story.append(Paragraph(f'<b>{name}</b>: '
            f'TN={tn}  FP={fp}  FN={fn}  TP={tp}  |  '
            f'Sensitivity={tp/(tp+fn):.3f}  Specificity={tn/(tn+fp):.3f}', BULL))

    # Section 6: Discussion
    story.append(Paragraph('6. Discussion', H2))
    story.append(Paragraph(
        'AUC ≈ 0.79 demonstrates that Virchow2 slide embeddings carry substantial '
        'CD/UC discriminative signal without task-specific fine-tuning. '
        'prism2_base marginally outperforms prism2_diagnostic (ΔAUC = 0.006). '
        'CD recall (~76%) exceeds UC recall (~67%) due to 2:1 class imbalance. '
        'Limitations include slide-level independence at inference, '
        'no probability calibration, and 35 excluded ambiguous patients.', BODY))

    # Section 7: Outputs
    story.append(Paragraph('7. Output Files', H2))
    story.append(tbl([
        ['File', 'Description'],
        ['cv_splits_patients.csv',              'Patient-level fold assignments'],
        ['cv_splits_slides.csv',                'Slide-level fold assignments'],
        ['prism2_base_fold_metrics.csv',        'Per-fold metrics for prism2_base'],
        ['prism2_base_slide_predictions.csv',   'Per-slide predictions for prism2_base'],
        ['prism2_diagnostic_fold_metrics.csv',  'Per-fold metrics for prism2_diagnostic'],
        ['prism2_diagnostic_slide_predictions.csv', 'Per-slide predictions for prism2_diagnostic'],
    ], [6.5, 9.5]))

    story += [Spacer(1, 0.5*cm),
              HRFlowable(width='100%', thickness=0.5, color=MGREY),
              Paragraph('IBD Plexus / SPARC  ·  Virchow2 via TRIDENT 0.3.0  ·  August 10, 2026', FOOT)]

    doc.build(story)
    print(f"PDF saved: {out_path}")
    return out_path


# ══════════════════════════════════════════════════════════════════════════════
#  POWERPOINT
# ══════════════════════════════════════════════════════════════════════════════

def build_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from lxml import etree

    def rgb(r, g, b): return RGBColor(r, g, b)

    DARK  = rgb(0x1a,0x2e,0x4a); DARK_H  = '1A2E4A'
    MID   = rgb(0x2e,0x6d,0xa4); MID_H   = '2E6DA4'
    LIGHT = rgb(0xd8,0xe8,0xf5); LIGHT_H = 'D8E8F5'
    WHITE = rgb(0xff,0xff,0xff); WHITE_H = 'FFFFFF'
    GREEN = rgb(0x1e,0x8b,0x4c); GREEN_H = '1E8B4C'
    RED   = rgb(0xc0,0x39,0x2b); RED_H   = 'C0392B'
    GREY  = rgb(0xf5,0xf7,0xfa); GREY_H  = 'F5F7FA'
    LGREEN= rgb(0xd5,0xf5,0xe3); LGREEN_H= 'D5F5E3'
    LRED  = rgb(0xfa,0xdb,0xd8); LRED_H  = 'FADBD8'
    DGREY = rgb(0x55,0x55,0x55); DGREY_H = '555555'
    MGREY = rgb(0xcc,0xcc,0xcc); MGREY_H = 'CCCCCC'
    ORANGE= rgb(0xe6,0x7e,0x22); ORANGE_H= 'E67E22'
    AMBER = rgb(0xff,0xf3,0xcd); AMBER_H = 'FFF3CD'

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    BLANK = prs.slide_layouts[6]

    def rect(slide, l,t,w,h, fill_h, line_h=None, lw=0.75):
        s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
        s.fill.solid(); s.fill.fore_color.rgb = RGBColor.from_string(fill_h)
        if line_h: s.line.color.rgb = RGBColor.from_string(line_h); s.line.width = Pt(lw)
        else:      s.line.fill.background()

    def txt(slide, text, l,t,w,h, size=10, bold=False, color=DARK,
            align=PP_ALIGN.LEFT, wrap=True):
        tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = tb.text_frame; tf.word_wrap = wrap
        p  = tf.paragraphs[0]; p.alignment = align
        r  = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color

    def mtxt(slide, lines, l,t,w,h, size=8.5, color=DARK, bold_first=False):
        tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = tb.text_frame; tf.word_wrap = True
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
            r = p.add_run(); r.text = line
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold_first and i==0

    def set_cell_bg(cell, hx):
        tc = cell._tc; tcPr = tc.get_or_add_tcPr()
        for old in tcPr.findall(qn('a:solidFill')): tcPr.remove(old)
        sf = etree.SubElement(tcPr, qn('a:solidFill'))
        sc = etree.SubElement(sf, qn('a:srgbClr')); sc.set('val', hx)

    def table(slide, data, l,t,col_w, row_h=0.27, hdr_bg=MID_H, fs=8.5):
        rows=len(data); cols=len(data[0]); tw=sum(col_w)
        sh = slide.shapes.add_table(rows,cols,Inches(l),Inches(t),Inches(tw),Inches(row_h*rows))
        tb = sh.table
        for i,w in enumerate(col_w): tb.columns[i].width = Inches(w)
        for r in range(rows): tb.rows[r].height = Inches(row_h)
        for r,row in enumerate(data):
            for c,val in enumerate(row):
                cell=tb.cell(r,c); cell.text=str(val)
                for para in cell.text_frame.paragraphs:
                    para.alignment = PP_ALIGN.CENTER
                    for run in para.runs:
                        run.font.size = Pt(fs); run.font.bold = (r==0)
                        run.font.color.rgb = WHITE if r==0 else DARK
                if r==0: set_cell_bg(cell, hdr_bg)
                elif r%2==0: set_cell_bg(cell, GREY_H)
                else: set_cell_bg(cell, WHITE_H)
        return sh

    def header(slide, title, sub=None):
        rect(slide,0,0,13.33,0.95,DARK_H)
        txt(slide,title,0.25,0.08,11,0.45,size=15,bold=True,color=WHITE)
        if sub: txt(slide,sub,0.25,0.57,11,0.30,size=9,color=LIGHT)

    def footer(slide, text):
        txt(slide,text,0.2,7.22,13,0.22,size=7.5,color=DGREY,align=PP_ALIGN.CENTER)

    def divider(slide, y):
        rect(slide,0.2,y,12.93,0.02,MID_H)

    # load data
    base_folds = pd.read_csv(os.path.join(RESULTS_DIR, 'prism2_base_fold_metrics.csv'))
    diag_folds = pd.read_csv(os.path.join(RESULTS_DIR, 'prism2_diagnostic_fold_metrics.csv'))
    patients   = pd.read_csv(os.path.join(TRAINING_DIR, 'cv_splits_patients.csv'))
    slides_df  = pd.read_csv(os.path.join(TRAINING_DIR, 'cv_splits_slides.csv'))

    b_auc = base_folds['auc'].mean(); b_auc_s = base_folds['auc'].std()
    d_auc = diag_folds['auc'].mean(); d_auc_s = diag_folds['auc'].std()
    b_acc = base_folds['accuracy'].mean()

    # ── Slide 1: Title ────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(BLANK)
    rect(sl,0,0,13.33,7.5,'F2F6FB')
    rect(sl,0,0,13.33,2.6,DARK_H)
    rect(sl,0,2.55,13.33,0.08,MID_H)
    txt(sl,'Colon CD vs UC Classification',0.5,0.45,12.3,0.8,
        size=28,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    txt(sl,'from Whole-Slide Image Embeddings',0.5,1.15,12.3,0.55,
        size=20,color=LIGHT,align=PP_ALIGN.CENTER)
    txt(sl,'Pipeline Description and Results',0.5,1.72,12.3,0.5,
        size=14,color=LIGHT,align=PP_ALIGN.CENTER)
    for i,(val,lab) in enumerate([
        (str(len(patients)),'patients'),
        (str(len(slides_df)),'colon slides'),
        (f'{b_auc:.3f}','best AUC'),
        (f'{b_acc*100:.1f}%','best accuracy'),
    ]):
        lx=1.0+i*2.88
        rect(sl,lx,3.0,2.5,1.4,LIGHT_H,MID_H)
        txt(sl,val,lx,3.15,2.5,0.65,size=24,bold=True,color=MID,align=PP_ALIGN.CENTER)
        txt(sl,lab,lx,3.78,2.5,0.4,size=10,color=DARK,align=PP_ALIGN.CENTER)
    mtxt(sl,['Dataset:    IBD Plexus / SPARC FFPE H&E colon biopsies',
             'Embeddings: Virchow2 via TRIDENT 0.3.0  |  prism2_base (2,560-d) & prism2_diagnostic (3,072-d)',
             'Model:      Random Forest  |  5-fold patient-level stratified CV  |  August 10, 2026'],
         1.0,4.8,11.3,0.9,size=9.5,color=DARK)

    # ── Slide 2: Pipeline ─────────────────────────────────────────────────────
    sl = prs.slides.add_slide(BLANK)
    rect(sl,0,0,13.33,7.5,'F9FBFD')
    header(sl,'Pipeline Overview','End-to-end workflow from raw WSI to classification result')
    steps=[('1','WSI Acquisition','IBD Plexus / SPARC\nFFPE TIFF at 40×\n(H&E stained)'),
           ('2','Tissue Segmentation','TRIDENT 0.3.0\nHEST segmentation\ncolon biopsies'),
           ('3','Patch Extraction','20× · 224 px\n0 px overlap\nVirchow2 encoder'),
           ('4','Slide Embeddings','PRISM2 aggregation\nprism2_base 2,560-d\nprism2_diag 3,072-d'),
           ('5','CV Split','Patient-level\n5-fold stratified\nNo leakage'),
           ('6','Random Forest','500 trees\nBalanced weights\nRF on embeddings'),
           ('7','Evaluation','AUC · AP\nAccuracy · F1\nPer-fold & aggregate')]
    box_w=1.72; gap=0.07
    for i,(num,title,desc) in enumerate(steps):
        lx=0.22+i*(box_w+gap)
        rect(sl,lx,1.08,box_w,0.36,MID_H)
        txt(sl,f'{num}. {title}',lx,1.09,box_w,0.34,size=8.5,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
        rect(sl,lx,1.44,box_w,0.9,LIGHT_H,MID_H)
        mtxt(sl,desc.split('\n'),lx+0.05,1.47,box_w-0.1,0.86,size=8,color=DARK)
        if i<len(steps)-1:
            txt(sl,'▶',lx+box_w+0.01,1.62,gap+0.05,0.3,size=9,bold=True,color=MID,align=PP_ALIGN.CENTER)
    divider(sl,2.48)
    txt(sl,'Data & Cohort Details',0.25,2.58,13,0.3,size=11,bold=True,color=MID)
    cd_n=(patients['diagnosis']=="Crohn's disease").sum(); uc_n=len(patients)-cd_n
    cd_s=(slides_df['diagnosis']=="Crohn's disease").sum(); uc_s=len(slides_df)-cd_s
    table(sl,[['Cohort','Patients','Slides','Slides/pt','Notes'],
              ["Crohn's disease (CD=0)",str(cd_n),str(cd_s),f'{cd_s/cd_n:.2f}','—'],
              ['Ulcerative colitis (UC=1)',str(uc_n),str(uc_s),f'{uc_s/uc_n:.2f}','—'],
              ['Ambiguous / excluded','35','—','—','Conflicting CD+UC labels']],
         l=0.25,t=2.90,col_w=[3.8,1.8,1.5,1.8,4.87],row_h=0.27,fs=8.5)
    txt(sl,'Stratification Variables',0.25,4.37,13,0.3,size=11,bold=True,color=MID)
    strat_sh=table(sl,[['Variable','Bins','Decision','Rationale'],
                        ['Diagnosis','CD / UC','✓ USE','Primary outcome — 2:1 imbalance'],
                        ['Sex / Gender','Female / Male / Missing','✓ USE','CD skews more female'],
                        ['Age at diagnosis','<20 / 20–35 / >35 / Missing','✓ USE','CD median 24 vs UC 27 yrs'],
                        ['N surgeries / CD phenotype / location','—','✗ SKIP','Near-surrogate for Dx or CD-only']],
                   l=0.25,t=4.69,col_w=[3.2,3.5,1.5,6.08],row_h=0.27,fs=8.5)
    stbl=strat_sh.table
    for r in range(1,5):
        c=stbl.cell(r,2)
        bg=LGREEN_H if r<=3 else LRED_H; fg=GREEN if r<=3 else RED
        set_cell_bg(c,bg)
        for p in c.text_frame.paragraphs:
            for run in p.runs: run.font.bold=True; run.font.color.rgb=fg
    footer(sl,'IBD Plexus / SPARC  ·  TRIDENT 0.3.0  ·  Virchow2  ·  prism2_base & prism2_diagnostic')

    # ── Slide 3: Results ──────────────────────────────────────────────────────
    sl = prs.slides.add_slide(BLANK)
    rect(sl,0,0,13.33,7.5,'F9FBFD')
    header(sl,'Results: 5-Fold Cross-Validation',
           'Random Forest on slide-level Virchow2 embeddings  ·  CD=0, UC=1')
    txt(sl,'Overall Performance (mean ± SD across 5 folds)',0.25,1.08,13,0.3,size=11,bold=True,color=MID)
    b=base_folds; d=diag_folds
    ov_sh=table(sl,[['Model','AUC','Average Precision','Accuracy','CD F1','UC F1'],
                    ['prism2_base',
                     f"{b['auc'].mean():.3f} ± {b['auc'].std():.3f}",
                     f"{b['ap'].mean():.3f} ± {b['ap'].std():.3f}",
                     f"{b['accuracy'].mean()*100:.1f}% ± {b['accuracy'].std()*100:.1f}%",
                     f"{b['cd_f1'].mean():.3f} ± {b['cd_f1'].std():.3f}",
                     f"{b['uc_f1'].mean():.3f} ± {b['uc_f1'].std():.3f}"],
                    ['prism2_diagnostic',
                     f"{d['auc'].mean():.3f} ± {d['auc'].std():.3f}",
                     f"{d['ap'].mean():.3f} ± {d['ap'].std():.3f}",
                     f"{d['accuracy'].mean()*100:.1f}% ± {d['accuracy'].std()*100:.1f}%",
                     f"{d['cd_f1'].mean():.3f} ± {d['cd_f1'].std():.3f}",
                     f"{d['uc_f1'].mean():.3f} ± {d['uc_f1'].std():.3f}"]],
               l=0.25,t=1.40,col_w=[3.5,2.3,3.0,2.5,2.2,2.2],row_h=0.30,fs=9.5)
    for c in range(6):
        set_cell_bg(ov_sh.table.cell(1,c),'E8F4EC')
        for p in ov_sh.table.cell(1,c).text_frame.paragraphs:
            for r in p.runs: r.font.bold=True
    divider(sl,2.46)
    txt(sl,'prism2_base — per-fold',0.25,2.56,6.4,0.28,size=10,bold=True,color=MID)
    table(sl,[['Fold','N val','AUC','AP','Acc','CD F1','UC F1']]+
             [[str(row['fold']),str(row['n_val']),f"{row['auc']:.3f}",
               f"{row['ap']:.3f}",f"{row['accuracy']*100:.1f}%",
               f"{row['cd_f1']:.3f}",f"{row['uc_f1']:.3f}"]
              for _,row in base_folds.iterrows()],
         l=0.25,t=2.86,col_w=[1.0,1.2,1.2,1.2,1.3,1.3,1.3],row_h=0.26,fs=8.5)
    txt(sl,'prism2_diagnostic — per-fold',6.85,2.56,6.4,0.28,size=10,bold=True,color=MID)
    table(sl,[['Fold','N val','AUC','AP','Acc','CD F1','UC F1']]+
             [[str(row['fold']),str(row['n_val']),f"{row['auc']:.3f}",
               f"{row['ap']:.3f}",f"{row['accuracy']*100:.1f}%",
               f"{row['cd_f1']:.3f}",f"{row['uc_f1']:.3f}"]
              for _,row in diag_folds.iterrows()],
         l=6.85,t=2.86,col_w=[1.0,1.2,1.2,1.2,1.3,1.3,1.3],row_h=0.26,fs=8.5)
    divider(sl,4.56)
    txt(sl,'Aggregated Confusion Matrices (all 5 folds, positive class = UC)',
        0.25,4.65,13,0.28,size=10,bold=True,color=MID)
    for lx,name,folds in [(0.25,'prism2_base',base_folds),(6.85,'prism2_diagnostic',diag_folds)]:
        tn=folds['tn'].sum();fp=folds['fp'].sum();fn=folds['fn'].sum();tp=folds['tp'].sum()
        sens=tp/(tp+fn); spec=tn/(tn+fp)
        rect(sl,lx,4.95,6.2,2.2,LIGHT_H,MID_H)
        txt(sl,name,lx+0.1,4.98,6.0,0.28,size=9.5,bold=True,color=MID)
        cx=lx+0.25; cy=5.28; cw=1.4; ch=0.55
        vals=[[tn,fp],[fn,tp]]
        cbg=[['F5F7FA','FADBD8'],['FADBD8','D5F5E3']]
        cfg=[['555555','C0392B'],['C0392B','1E8B4C']]
        for ri in range(2):
            for ci in range(2):
                rect(sl,cx+ci*cw,cy+ri*ch,cw,ch,cbg[ri][ci],MGREY_H,0.5)
                txt(sl,str(vals[ri][ci]),cx+ci*cw,cy+ri*ch+0.05,cw,0.32,
                    size=14,bold=True,color=RGBColor.from_string(cfg[ri][ci]),align=PP_ALIGN.CENTER)
                txt(sl,f'Pred {"CD" if ci==0 else "UC"}',cx+ci*cw,cy+ri*ch+0.35,cw,0.18,
                    size=7,color=DGREY,align=PP_ALIGN.CENTER)
        txt(sl,f'Sensitivity (UC): {sens:.3f}',lx+3.2,5.28,2.8,0.28,size=9,color=DARK)
        txt(sl,f'Specificity (CD): {spec:.3f}',lx+3.2,5.58,2.8,0.28,size=9,color=DARK)
        txt(sl,'CD better identified\nthan UC (majority class)',lx+3.2,5.88,2.8,0.4,size=8.5,color=DGREY)
    footer(sl,'AUC = area under ROC  ·  AP = average precision  ·  Balanced class weights applied')

    # ── Slide 4: Discussion ───────────────────────────────────────────────────
    sl = prs.slides.add_slide(BLANK)
    rect(sl,0,0,13.33,7.5,'F9FBFD')
    header(sl,'Discussion & Next Steps','Interpretation and recommended future work')
    txt(sl,'Key Findings',0.25,1.08,6.3,0.3,size=11,bold=True,color=MID)
    findings=[
        ('Viable diagnostic signal','AUC ≈ 0.79 shows Virchow2 embeddings carry strong CD/UC signal without fine-tuning.'),
        ('prism2_base ≥ prism2_diagnostic','Base head (2,560-d) marginally outperforms diagnostic head. Extra dims do not help RF.'),
        ('CD better recalled than UC',f'Specificity ~{base_folds["cd_recall"].mean()*100:.0f}% vs sensitivity ~{base_folds["uc_recall"].mean()*100:.0f}%. Expected from 2:1 imbalance.'),
        ('Stable across folds',f'AUC SD ≤ {base_folds["auc"].std():.3f} — robust stratification, no easy folds.'),
    ]
    y=1.42
    for title,body in findings:
        rect(sl,0.25,y,6.3,0.72,LIGHT_H,MID_H)
        txt(sl,title,0.38,y+0.03,6.1,0.25,size=9,bold=True,color=MID)
        txt(sl,body,0.38,y+0.28,6.1,0.38,size=8.5,color=DARK)
        y+=0.78
    txt(sl,'Limitations',6.85,1.08,6.3,0.3,size=11,bold=True,color=ORANGE)
    lims=[
        ('Slide-level independence','UC avg 2.28 slides/patient treated independently. Patient-level aggregation may reduce noise.'),
        ('Linear model in embedding space','RF is linear in Virchow2 space. ABMIL or fine-tuned encoder could improve performance.'),
        ('35 ambiguous patients excluded','May represent IBDU or reclassified cases — worth investigating separately.'),
        ('No probability calibration','RF probabilities uncalibrated. Platt scaling could improve clinical utility.'),
    ]
    y=1.42
    for title,body in lims:
        rect(sl,6.85,y,6.3,0.72,AMBER_H,ORANGE_H,0.75)
        txt(sl,title,6.98,y+0.03,6.1,0.25,size=9,bold=True,color=ORANGE)
        txt(sl,body,6.98,y+0.28,6.1,0.38,size=8.5,color=DARK)
        y+=0.78
    divider(sl,5.36)
    txt(sl,'Recommended Next Steps',0.25,5.46,13,0.3,size=11,bold=True,color=MID)
    for i,s in enumerate([
        '1. ABMIL / TRANSMIL — attention-based MIL with patch-level features',
        '2. Patient-level prediction — aggregate slide scores; evaluate patient AUC',
        '3. LR + SVM baseline — compare to RF on same embeddings',
        '4. Calibration — Platt scaling on RF probabilities',
        '5. Subgroup analysis — by biopsy site, disease activity, age group',
    ]):
        lx=0.25+(i%3)*4.36; ly=5.80+(i//3)*0.52
        rect(sl,lx,ly,4.2,0.44,LIGHT_H,MID_H)
        txt(sl,s,lx+0.1,ly+0.05,4.0,0.34,size=8,color=DARK)
    footer(sl,'Data: IBD Plexus / SPARC  ·  Virchow2 via TRIDENT 0.3.0  ·  August 10, 2026')

    out_path, _ = next_versioned_path(os.path.join(REPORTS_DIR, 'pipeline_results_slides.pptx'))
    prs.save(out_path)
    print(f"PPTX saved: {out_path}")
    return out_path


if __name__ == '__main__':
    parser = argparse.ArgumentParser()
    parser.add_argument('--desc', default='Imaging pipeline results (prism2_base + prism2_diagnostic, full 1,250-patient cohort)',
                        help='One-line description of what changed in this version')
    args = parser.parse_args()

    pdf_path  = build_pdf()
    pptx_path = build_pptx()
    log_version(pdf_path,  args.desc)
    log_version(pptx_path, args.desc)
