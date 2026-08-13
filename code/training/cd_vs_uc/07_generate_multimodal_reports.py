"""
Generate PDF report and PowerPoint slides for the multimodal fusion study.

Reads pre-computed fold-metric CSVs and summary JSONs; does NOT re-run any model.

Outputs  (versioned)
--------
  multimodal_report_vN.pdf
  multimodal_slides_vN.pptx
  VERSIONS.md  (append-only changelog)

Usage
-----
  python 07_generate_multimodal_reports.py
  python 07_generate_multimodal_reports.py --desc "Added site-matched results"
"""

import argparse
import json
import os
import numpy as np
import pandas as pd
from version_utils import next_versioned_path, log_version

MM_DIR      = '/home/jovyan/kgbk271-ibd-volume/training/cd_vs_uc/06_07_multimodal_allsites/results'
REPORTS_DIR = '/home/jovyan/kgbk271-ibd-volume/training/cd_vs_uc/06_07_multimodal_allsites/reports'
TRANS_DIR   = '/home/jovyan/kgbk271-ibd-volume/training/cd_vs_uc/03_05_transcriptomics_allsites/results'
IMG_DIR     = '/home/jovyan/kgbk271-ibd-volume/training/cd_vs_uc/02_04_imaging_allsites/results'


# ── data loading ──────────────────────────────────────────────────────────────

def load_all():
    mm   = pd.read_csv(os.path.join(MM_DIR, 'multimodal_fold_metrics.csv'))
    abl  = pd.read_csv(os.path.join(MM_DIR, 'multimodal_ablation_fold_metrics.csv'))
    rna  = pd.read_csv(os.path.join(TRANS_DIR, 'transcriptomics_vst_fold_metrics.csv'))
    bm   = pd.read_csv(os.path.join(IMG_DIR, 'prism2_base_matched_fold_metrics.csv'))
    return mm, abl, rna, bm


def strat(df, name):
    """Return fold-metrics rows for one strategy."""
    return df[df['strategy'] == name]


def summ(rows):
    a = rows['auc'].values; p = rows['ap'].values; c = rows['accuracy'].values
    return (round(float(a.mean()),4), round(float(a.std()),4),
            round(float(p.mean()),4), round(float(p.std()),4),
            round(float(c.mean()),4), round(float(c.std()),4))


def cm_stats(rows):
    tn = rows['tn'].sum(); fp = rows['fp'].sum()
    fn = rows['fn'].sum(); tp = rows['tp'].sum()
    return tn, fp, fn, tp, tp/(tp+fn), tn/(tn+fp)


# ── strategy metadata ─────────────────────────────────────────────────────────

STRATEGIES = [
    # (display_name, csv_name, source_df_key, dim, category)
    ('img_base (full)',           'img_base_patmean',        'mm',  '2,560',    'imaging'),
    ('img_diag (full)',           'img_diag_patmean',        'mm',  '3,072',    'imaging'),
    ('RNA VST (full)',            'rna_patmean',             'mm',  '17,963',   'rna'),
    ('late fusion (avg)',         'late_fusion_base_rna',    'mm',  '2+17k',    'fusion'),
    ('concat raw / scaled',       'concat_scaled_base_rna',  'mm',  '20,523',   'fusion'),
    ('concat PCA-128',            'concat_pca128_base_rna',  'mm',  '2×128',    'fusion'),
    ('site-matched scaled',       'site_matched_scaled_base','mm',  '20,523',   'site'),
    ('img_base PCA-128 (uni)',    'img_base_pca128',         'abl', '128',      'ablation'),
    ('RNA PCA-128 (uni)',         'rna_pca128',              'abl', '128',      'ablation'),
    ('concat raw (no scale)',     'concat_raw',              'abl', '20,523',   'ablation'),
]


# ══════════════════════════════════════════════════════════════════════════════
#  PDF REPORT
# ══════════════════════════════════════════════════════════════════════════════

def build_pdf():
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.lib import colors
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer, Table,
                                    TableStyle, HRFlowable, PageBreak)
    from reportlab.lib.enums import TA_JUSTIFY, TA_CENTER

    DARK   = colors.HexColor('#1a2e4a')
    MID    = colors.HexColor('#2e6da4')
    TEAL   = colors.HexColor('#16a085')
    PURPLE = colors.HexColor('#6c3082')
    LTEAL  = colors.HexColor('#d1f0eb')
    LPURP  = colors.HexColor('#e8d5f5')
    LIGHT  = colors.HexColor('#d8e8f5')
    WHITE  = colors.white
    GREY   = colors.HexColor('#f5f7fa')
    MGREY  = colors.HexColor('#dddddd')
    DGREY  = colors.HexColor('#888888')
    ORANGE = colors.HexColor('#e67e22')
    GREEN  = colors.HexColor('#1e8b4c')
    LGREEN = colors.HexColor('#d5f5e3')
    RED    = colors.HexColor('#c0392b')
    AMBER  = colors.HexColor('#fff3cd')

    mm, abl, rna_f, bm_f = load_all()
    dfs = {'mm': mm, 'abl': abl}

    out_path, _ = next_versioned_path(os.path.join(REPORTS_DIR, 'multimodal_report.pdf'))
    doc = SimpleDocTemplate(out_path, pagesize=A4,
                            topMargin=2.2*cm, bottomMargin=2*cm,
                            leftMargin=2.2*cm, rightMargin=2.2*cm)

    S = getSampleStyleSheet()
    def sty(name, **kw): return ParagraphStyle(name, parent=S['Normal'], **kw)

    H1   = sty('H1',   fontSize=18, textColor=DARK,   spaceAfter=6,  fontName='Helvetica-Bold')
    H2   = sty('H2',   fontSize=13, textColor=PURPLE,  spaceBefore=16, spaceAfter=5, fontName='Helvetica-Bold')
    H3   = sty('H3',   fontSize=11, textColor=DARK,   spaceBefore=10, spaceAfter=3, fontName='Helvetica-Bold')
    BODY = sty('BODY', fontSize=9.5, leading=15, spaceAfter=4, alignment=TA_JUSTIFY)
    SMALL= sty('SMALL',fontSize=8.5, leading=12, textColor=DGREY)
    BULL = sty('BULL', fontSize=9.5, leading=14, leftIndent=14, spaceAfter=3)
    NOTE = sty('NOTE', fontSize=8.5, leading=12, textColor=colors.HexColor('#444444'), leftIndent=10)
    FOOT = sty('FOOT', fontSize=7.5, textColor=DGREY, alignment=TA_CENTER)

    def TS(hdr=MID): return [
        ('BACKGROUND', (0,0), (-1,0), hdr), ('TEXTCOLOR', (0,0), (-1,0), WHITE),
        ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'), ('FONTSIZE', (0,0), (-1,-1), 9),
        ('ROWBACKGROUNDS', (0,1), (-1,-1), [WHITE, GREY]),
        ('GRID', (0,0), (-1,-1), 0.4, MGREY), ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
        ('TOPPADDING', (0,0), (-1,-1), 4), ('BOTTOMPADDING', (0,0), (-1,-1), 4),
        ('LEFTPADDING', (0,0), (-1,-1), 6), ('RIGHTPADDING', (0,0), (-1,-1), 6),
    ]
    def tbl(data, cw, hdr=MID, extra=None):
        t = Table(data, colWidths=[w*cm for w in cw])
        t.setStyle(TableStyle(TS(hdr) + (extra or [])))
        return t

    rna_auc, _, _, _, _, _ = summ(rna_f)
    bm_auc,  _, _, _, _, _ = summ(bm_f)
    best_fusion = summ(strat(mm, 'concat_scaled_base_rna'))
    img_base_s  = summ(strat(mm, 'img_base_patmean'))
    rna_s       = summ(strat(mm, 'rna_patmean'))

    story = []

    # ── Title ──────────────────────────────────────────────────────────────────
    story += [
        Spacer(1, 1.5*cm),
        Paragraph('Multimodal Fusion: Imaging + Transcriptomics', H1),
        Paragraph('CD vs UC Classification · Patient-Level Feature Fusion · 997 Matched Patients',
                  sty('sub', fontSize=13, textColor=PURPLE, spaceAfter=6)),
        HRFlowable(width='100%', thickness=2, color=PURPLE, spaceAfter=10),
        Paragraph('Dataset: IBD Plexus / SPARC  ·  Imaging: Virchow2 prism2_base (2,560-d)  ·  '
                  'RNA: CombatSeq VST (17,963 genes)  ·  August 11, 2026', SMALL),
        Spacer(1, 0.5*cm),
    ]
    story.append(tbl([
        ['Best fusion AUC', f"{best_fusion[0]:.3f} ± {best_fusion[1]:.3f}  (concat_raw / concat_scaled)"],
        ['RNA only AUC',    f"{rna_s[0]:.3f} ± {rna_s[1]:.3f}  (patient-mean, 17,963 features)"],
        ['Imaging only AUC',f"{img_base_s[0]:.3f} ± {img_base_s[1]:.3f}  (patient-mean, 2,560 features)"],
        ['ΔAUC fusion vs RNA', f"{best_fusion[0]-rna_s[0]:+.4f}  (fusion adds minimal signal over RNA alone)"],
        ['Matched cohort', '997 patients with both colon WSI and colon RNAseq · identical 5-fold CV'],
    ], [4, 12], hdr=PURPLE, extra=[
        ('FONTNAME', (0,0), (0,-1), 'Helvetica-Bold'),
        ('BACKGROUND', (0,0), (0,-1), LPURP),
        ('BACKGROUND', (0,0), (-1,0), LPURP),
        ('TEXTCOLOR', (0,0), (-1,0), DARK),
    ]))
    story.append(PageBreak())

    # ── Section 1: Design ──────────────────────────────────────────────────────
    story.append(Paragraph('1. Study Design', H2))
    story.append(Paragraph('1.1  Patient-level aggregation', H3))
    story.append(Paragraph(
        'Imaging and RNA are at different granularities: a patient may have multiple '
        'colon slides and multiple RNA samples from different biopsy sites. '
        'To create a single feature vector per patient for each modality, all slides '
        '(or RNA samples) from that patient are <b>mean-pooled</b> before fusion. '
        'This produces a stable whole-colon representation and ensures a 1:1 patient ratio '
        'regardless of sampling depth.', BODY))

    story.append(Paragraph('1.2  Biopsy-site mapping validation', H3))
    story.append(Paragraph(
        'Before mean-pooling, we confirmed the RNA ↔ imaging anatomical correspondence '
        'by matching samples at the (patient, location) level:', BODY))
    story.append(tbl([
        ['Mapping check', 'Count'],
        ['Patients with both RNA and imaging',         '1,007'],
        ['Exact location match (same site set)',        '856 (85%)'],
        ['Partial overlap (some sites match)',          '137'],
        ['No overlap (different sites entirely)',        '14'],
        ['Patient-location pairs with BOTH modalities','1,388  (993 patients)'],
        ['RNA-only patient-location pairs',             '71'],
        ['Imaging-only patient-location pairs',         '441'],
    ], [9, 7], hdr=PURPLE, extra=[
        ('BACKGROUND', (0,5), (-1,5), LPURP),
        ('FONTNAME', (0,5), (-1,5), 'Helvetica-Bold'),
    ]))
    story.append(Paragraph(
        'Conclusion: 85% of patients have exact RNA ↔ imaging location correspondence. '
        'The 14 patients with no overlap are retained in the patient-mean fusion '
        '(their RNA and imaging come from different colon segments, but the CD/UC label '
        'is patient-level and site-independent).', NOTE))

    # ── Section 2: Fusion Strategies ──────────────────────────────────────────
    story.append(Paragraph('2. Fusion Strategies', H2))
    story.append(tbl([
        ['Strategy', 'Modalities', 'Preprocessing', 'Dim', 'Category'],
        ['img_base_patmean',         'Imaging',       'Mean pool',              '2,560',  'Unimodal baseline'],
        ['rna_patmean',              'RNA',           'Mean pool',              '17,963', 'Unimodal baseline'],
        ['late_fusion',              'Img + RNA',     'Separate RFs, avg P(UC)','2+17k',  'Score fusion'],
        ['concat_raw',               'Img + RNA',     'None (raw concat)',       '20,523', 'Feature fusion'],
        ['concat_scaled',            'Img + RNA',     'StandardScaler per block','20,523', 'Feature fusion'],
        ['concat_pca128',            'Img + RNA',     'PCA-128 per block',       '256',    'Feature fusion'],
        ['site_matched_scaled',      'Img + RNA',     'Scaled, site-matched pairs','20,523','Site-matched'],
        ['img_base_pca128 (abl.)',   'Imaging',       'PCA-128 projection',      '128',    'PCA ablation'],
        ['rna_pca128 (abl.)',        'RNA',           'PCA-128 projection',      '128',    'PCA ablation'],
    ], [3.8, 2.2, 4.0, 2.0, 4.0], hdr=PURPLE))

    # ── Section 3: Results ────────────────────────────────────────────────────
    story.append(Paragraph('3. Results', H2))
    story.append(Paragraph('3.1  Main fusion comparison', H3))

    rows = [['Strategy', 'Dim', 'AUC', 'AP', 'Accuracy', 'CD F1', 'UC F1']]
    for display, csv_name, src, dim, cat in STRATEGIES:
        df = dfs[src]
        r = strat(df, csv_name)
        if len(r) == 0: continue
        a, as_, p, ps, c, cs = summ(r)
        cd_f1 = r['cd_f1'].mean(); uc_f1 = r['uc_f1'].mean()
        rows.append([display, dim,
                     f'{a:.3f}±{as_:.3f}', f'{p:.3f}±{ps:.3f}',
                     f'{c*100:.1f}%±{cs*100:.1f}%',
                     f'{cd_f1:.3f}', f'{uc_f1:.3f}'])
    story.append(tbl(rows, [4.0, 2.2, 2.8, 2.8, 2.8, 1.7, 1.7], hdr=PURPLE, extra=[
        ('BACKGROUND', (0,3), (-1,3), LPURP),
        ('FONTNAME',   (0,3), (-1,3), 'Helvetica-Bold'),
        ('BACKGROUND', (0,4), (-1,4), LPURP),
        ('FONTNAME',   (0,4), (-1,4), 'Helvetica-Bold'),
    ]))
    story.append(Paragraph(
        'concat_raw and concat_scaled are highlighted — they achieve the best AUC and are '
        'identical because Random Forest is invariant to monotonic feature scaling '
        '(see Section 4.2).', NOTE))

    story.append(Paragraph('3.2  Site-matched vs patient-mean fusion', H3))
    sm_s = summ(strat(mm, 'site_matched_scaled_base'))
    cs_s = summ(strat(mm, 'concat_scaled_base_rna'))
    story.append(Paragraph(
        f'Site-matched fusion (AUC {sm_s[0]:.3f}) is substantially worse than patient-mean '
        f'fusion (AUC {cs_s[0]:.3f}). The reason: the CD/UC label is patient-level; '
        f'averaging across all colon sites reduces per-sample noise and yields a more '
        f'stable patient representation. Pairing individual biopsy sites adds '
        f'variance without adding label-relevant information.', BODY))

    # ── Section 4: Ablation ───────────────────────────────────────────────────
    story.append(Paragraph('4. Ablation Analysis', H2))
    story.append(Paragraph('4.1  Effect of PCA compression per modality', H3))

    img_full = summ(strat(mm,  'img_base_patmean'))
    img_pca  = summ(strat(abl, 'img_base_pca128'))
    rna_full = summ(strat(mm,  'rna_patmean'))
    rna_pca  = summ(strat(abl, 'rna_pca128'))
    con_pca  = summ(strat(mm,  'concat_pca128_base_rna'))
    con_raw  = summ(strat(abl, 'concat_raw'))

    pca_var_img = abl[abl['strategy']=='img_base_pca128']['pca_var_explained'].mean()
    pca_var_rna = abl[abl['strategy']=='rna_pca128']['pca_var_explained'].mean()

    story.append(tbl([
        ['Modality', 'Full dim AUC', 'PCA-128 AUC', 'ΔAUC', 'Variance retained'],
        ['Imaging (prism2_base)',
         f'{img_full[0]:.3f}', f'{img_pca[0]:.3f}',
         f'{img_pca[0]-img_full[0]:+.3f}', f'{pca_var_img*100:.1f}%'],
        ['RNA (VST)',
         f'{rna_full[0]:.3f}', f'{rna_pca[0]:.3f}',
         f'{rna_pca[0]-rna_full[0]:+.3f}', f'{pca_var_rna*100:.1f}%'],
        ['Concat (img + RNA)',
         f'{con_raw[0]:.3f}', f'{con_pca[0]:.3f}',
         f'{con_pca[0]-con_raw[0]:+.3f}', '—'],
    ], [4.5, 3, 3, 2.5, 3], hdr=PURPLE, extra=[
        ('BACKGROUND', (0,2), (-1,2), LPURP),
        ('FONTNAME',   (0,2), (-1,2), 'Helvetica-Bold'),
    ]))
    story.append(Paragraph(
        '<b>Key finding:</b> Imaging is nearly intrinsically low-dimensional — '
        f'128 PCs capture {pca_var_img*100:.1f}% of its variance, yet AUC still '
        f'drops by {abs(img_pca[0]-img_full[0]):.3f}. '
        f'RNA is high-dimensional — 128 PCs capture only {pca_var_rna*100:.1f}%, '
        f'and the AUC drops by {abs(rna_pca[0]-rna_full[0]):.3f}. '
        'The drop in concat_pca128 is therefore driven primarily by RNA compression, '
        'not imaging compression.', BODY))

    story.append(Paragraph('4.2  Raw vs scaled concatenation', H3))
    story.append(Paragraph(
        '<b>concat_raw and concat_scaled produce identical results fold-for-fold.</b> '
        'Random Forest uses rank-based splits (Gini impurity on thresholded values), '
        'which are invariant to monotonic feature transformations. '
        'StandardScaler changes the numeric scale but preserves feature ordering, '
        'so the RF makes the exact same decisions. '
        'The concern that RNA\'s 17,963 features would numerically dominate imaging\'s '
        '2,560 features via scale does <b>not apply to Random Forest</b>. '
        'Features dominate the model only if they are more predictive, not because '
        'of their numeric range.', BODY))

    # ── Section 5: Discussion ─────────────────────────────────────────────────
    story.append(Paragraph('5. Discussion', H2))
    story.append(Paragraph(
        f'Fusion provides no meaningful improvement over RNA alone (ΔAUC = '
        f'{best_fusion[0]-rna_full[0]:+.4f}). The RNA modality already captures '
        f'nearly all discriminative signal between CD and UC. '
        f'Imaging adds marginal complementary information at best. '
        f'This does not mean fusion is useless — the confidence of the combined '
        f'model may be better calibrated, and a non-linear fusion model '
        f'(e.g. neural network late fusion or cross-modal attention) may extract '
        f'interactions that a RF on concatenated features cannot.', BODY))
    story.append(Paragraph(
        'The site-matched analysis confirms that patient-mean pooling is the correct '
        'aggregation strategy when the prediction target is patient-level (CD vs UC). '
        'Site-level pairing increases variance without improving the label signal.', BODY))

    story.append(Paragraph('6. Next Steps', H2))
    for s in [
        '1. Neural late fusion — learned weighted combination of RNA and imaging scores',
        '2. Cross-modal attention — transformer that attends over imaging patches conditioned on RNA',
        '3. Site-level CD/UC features — test if site-specific expression patterns add signal over patient mean',
        '4. SHAP on concat model — identify which genes vs. imaging dimensions drive predictions',
        '5. Calibration — Platt scaling on RF probabilities before fusion',
        '6. Biopsy-site correction — strip site-covariate signal before fusion',
    ]:
        story.append(Paragraph(s, BULL))

    story += [
        Spacer(1, 0.5*cm),
        HRFlowable(width='100%', thickness=0.5, color=MGREY),
        Paragraph('IBD Plexus / SPARC  ·  Virchow2 TRIDENT 0.3.0  ·  '
                  'CombatSeq VST GSF1491805  ·  August 11, 2026', FOOT),
    ]

    doc.build(story)
    print(f'PDF saved: {out_path}')
    return out_path


# ══════════════════════════════════════════════════════════════════════════════
#  POWERPOINT
# ══════════════════════════════════════════════════════════════════════════════

def build_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from lxml import etree

    mm, abl, rna_f, bm_f = load_all()
    dfs = {'mm': mm, 'abl': abl}

    def rgb(r,g,b): return RGBColor(r,g,b)
    DARK   = rgb(0x1a,0x2e,0x4a); DARK_H   = '1A2E4A'
    MID    = rgb(0x2e,0x6d,0xa4); MID_H    = '2E6DA4'
    TEAL   = rgb(0x16,0xa0,0x85); TEAL_H   = '16A085'
    PURP   = rgb(0x6c,0x30,0x82); PURP_H   = '6C3082'
    LPURP  = rgb(0xe8,0xd5,0xf5); LPURP_H  = 'E8D5F5'
    LTEAL  = rgb(0xd1,0xf0,0xeb); LTEAL_H  = 'D1F0EB'
    LIGHT  = rgb(0xd8,0xe8,0xf5); LIGHT_H  = 'D8E8F5'
    WHITE  = rgb(0xff,0xff,0xff); WHITE_H  = 'FFFFFF'
    GREEN  = rgb(0x1e,0x8b,0x4c); GREEN_H  = '1E8B4C'
    LGREEN_H= 'D5F5E3'
    GREY   = rgb(0xf5,0xf7,0xfa); GREY_H   = 'F5F7FA'
    DGREY  = rgb(0x55,0x55,0x55); DGREY_H  = '555555'
    MGREY_H= 'CCCCCC'
    ORANGE = rgb(0xe6,0x7e,0x22); ORANGE_H = 'E67E22'
    AMBER_H= 'FFF3CD'
    RED    = rgb(0xc0,0x39,0x2b); RED_H    = 'C0392B'
    LRED_H = 'FADBD8'

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    BLANK = prs.slide_layouts[6]

    def rect(sl,l,t,w,h,fh,lh=None,lw=0.75):
        s = sl.shapes.add_shape(1,Inches(l),Inches(t),Inches(w),Inches(h))
        s.fill.solid(); s.fill.fore_color.rgb = RGBColor.from_string(fh)
        if lh: s.line.color.rgb = RGBColor.from_string(lh); s.line.width = Pt(lw)
        else:  s.line.fill.background()

    def txt(sl,text,l,t,w,h,size=10,bold=False,color=DARK,align=PP_ALIGN.LEFT,wrap=True):
        tb = sl.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
        tf = tb.text_frame; tf.word_wrap = wrap
        p = tf.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color

    def mtxt(sl,lines,l,t,w,h,size=8.5,color=DARK):
        tb = sl.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
        tf = tb.text_frame; tf.word_wrap = True
        for i,line in enumerate(lines):
            p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
            r = p.add_run(); r.text = line
            r.font.size = Pt(size); r.font.color.rgb = color

    def set_cell_bg(cell, hx):
        tc = cell._tc; tcPr = tc.get_or_add_tcPr()
        for old in tcPr.findall(qn('a:solidFill')): tcPr.remove(old)
        sf = etree.SubElement(tcPr, qn('a:solidFill'))
        sc = etree.SubElement(sf, qn('a:srgbClr')); sc.set('val', hx)

    def table(sl,data,l,t,cw,rh=0.27,hdr=PURP_H,fs=8.5):
        rows=len(data); cols=len(data[0]); tw=sum(cw)
        sh = sl.shapes.add_table(rows,cols,Inches(l),Inches(t),Inches(tw),Inches(rh*rows))
        tb = sh.table
        for i,w in enumerate(cw): tb.columns[i].width = Inches(w)
        for r in range(rows): tb.rows[r].height = Inches(rh)
        for r,row in enumerate(data):
            for c,val in enumerate(row):
                cell = tb.cell(r,c); cell.text = str(val)
                for para in cell.text_frame.paragraphs:
                    para.alignment = PP_ALIGN.CENTER
                    for run in para.runs:
                        run.font.size = Pt(fs); run.font.bold = (r==0)
                        run.font.color.rgb = WHITE if r==0 else DARK
                if r==0: set_cell_bg(cell, hdr)
                elif r%2==0: set_cell_bg(cell, GREY_H)
                else: set_cell_bg(cell, WHITE_H)
        return sh

    def header(sl,title,sub=None):
        rect(sl,0,0,13.33,0.95,DARK_H)
        txt(sl,title,0.25,0.08,11,0.45,size=15,bold=True,color=WHITE)
        if sub: txt(sl,sub,0.25,0.57,11,0.30,size=9,color=LIGHT)

    def footer(sl,text):
        txt(sl,text,0.2,7.22,13,0.22,size=7.5,color=DGREY,align=PP_ALIGN.CENTER)

    def divider(sl,y):
        rect(sl,0.2,y,12.93,0.02,PURP_H)

    # convenience
    def s(csv_name, src='mm'):
        r = strat(dfs[src], csv_name)
        return summ(r) if len(r) else None

    best_s   = s('concat_scaled_base_rna')
    rna_s    = s('rna_patmean')
    img_s    = s('img_base_patmean')
    img_pca_s= s('img_base_pca128', 'abl')
    rna_pca_s= s('rna_pca128', 'abl')
    raw_s    = s('concat_raw', 'abl')
    pca_c_s  = s('concat_pca128_base_rna')

    pca_var_img = abl[abl['strategy']=='img_base_pca128']['pca_var_explained'].mean()
    pca_var_rna = abl[abl['strategy']=='rna_pca128']['pca_var_explained'].mean()

    # ── Slide 1: Title ────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(BLANK)
    rect(sl,0,0,13.33,7.5,'F5F0FA')
    rect(sl,0,0,13.33,2.7,DARK_H)
    rect(sl,0,2.65,13.33,0.08,PURP_H)
    txt(sl,'Multimodal Fusion',0.5,0.32,12.3,0.68,
        size=28,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    txt(sl,'Imaging + Transcriptomics for CD vs UC Classification',0.5,0.98,12.3,0.5,
        size=18,color=LIGHT,align=PP_ALIGN.CENTER)
    txt(sl,'997 matched patients  ·  Patient-level mean pooling  ·  6 fusion strategies + ablations  ·  Random Forest',
        0.5,1.60,12.3,0.4,size=10.5,color=LIGHT,align=PP_ALIGN.CENTER)
    for i,(val,lab,fh,bh) in enumerate([
        (f'{best_s[0]:.3f}','best fusion AUC',    PURP_H,  LPURP_H),
        (f'{rna_s[0]:.3f}', 'RNA only AUC',       TEAL_H,  LTEAL_H),
        (f'{img_s[0]:.3f}', 'imaging only AUC',   MID_H,   LIGHT_H),
        (f'{best_s[0]-rna_s[0]:+.4f}','ΔAUC fuse vs RNA', PURP_H, LPURP_H),
    ]):
        lx = 1.0 + i*2.88
        rect(sl,lx,3.05,2.5,1.4,bh,fh)
        txt(sl,val,lx,3.18,2.5,0.65,size=24,bold=True,
            color=PURP if fh==PURP_H else (TEAL if fh==TEAL_H else MID),
            align=PP_ALIGN.CENTER)
        txt(sl,lab,lx,3.82,2.5,0.4,size=9.5,color=DARK,align=PP_ALIGN.CENTER)
    mtxt(sl,['Imaging: Virchow2 prism2_base (2,560-d) & prism2_diagnostic (3,072-d)',
             'Transcriptomics: CombatSeq VST  ·  17,963 genes  ·  997 patients mean-pooled',
             'Same patient-level 5-fold CV splits throughout  ·  August 11, 2026'],
         1.0,4.85,11.3,0.9,size=9.5,color=DARK)

    # ── Slide 2: Design + biopsy site validation ──────────────────────────────
    sl = prs.slides.add_slide(BLANK)
    rect(sl,0,0,13.33,7.5,'F9FBFD')
    header(sl,'Study Design','Patient-level aggregation · biopsy-site mapping validation')

    txt(sl,'Patient-Level Mean Pooling',0.25,1.08,6.3,0.30,size=11,bold=True,color=PURP)
    for i,(title,body) in enumerate([
        ('Why mean-pool?','A patient may have 3 slides from different colon sites and 2 RNA samples. Mean pooling gives one patient vector per modality — comparable across patients regardless of sampling depth.'),
        ('Unit of analysis','All 997 patients contribute exactly one imaging vector (2,560-d) and one RNA vector (17,963-d). Fusion operates at this patient level.'),
        ('Fold assignment','Patient-level folds prevent any patient appearing in both train and validation, even when multiple slides/samples exist.'),
    ]):
        ly = 1.42 + i*0.80
        rect(sl,0.25,ly,6.3,0.72,LPURP_H,PURP_H)
        txt(sl,title,0.38,ly+0.03,6.1,0.25,size=9,bold=True,color=PURP)
        txt(sl,body, 0.38,ly+0.28,6.1,0.40,size=8.5,color=DARK)

    divider(sl,3.85)
    txt(sl,'Biopsy-Site Mapping Validation',0.25,3.95,6.3,0.30,size=11,bold=True,color=PURP)
    sh=table(sl,[['Check','Count','Note'],
                 ['Patients with both RNA and imaging','1,007','pre-filter'],
                 ['Exact location match (same site set)','856 (85%)','✓'],
                 ['Partial overlap','137','✓ (some sites shared)'],
                 ['No site overlap at all','14','retained in patient mean'],
                 ['(Patient, location) pairs with BOTH','1,388 / 993 pts','used in site-matched exp.'],
                 ['RNA-only (patient, location) pairs','71','—'],
                 ['Imaging-only (patient, location) pairs','441','—']],
             l=0.25,t=4.30,cw=[5.0,2.6,1.5],rh=0.26,fs=8.5)
    for c in range(3): set_cell_bg(sh.table.cell(5,c), LPURP_H)

    txt(sl,'Fusion Strategies Overview',7.0,1.08,6.1,0.30,size=11,bold=True,color=PURP)
    sh2=table(sl,[['Strategy','Dim','Type'],
                  ['img_base_patmean','2,560','Unimodal'],
                  ['rna_patmean','17,963','Unimodal'],
                  ['late_fusion','2+17k','Score avg'],
                  ['concat_raw / scaled','20,523','Feature'],
                  ['concat_pca128','256','Feature+PCA'],
                  ['site_matched_scaled','20,523','Site-paired'],
                  ['img_base_pca128 (abl)','128','PCA ablation'],
                  ['rna_pca128 (abl)','128','PCA ablation']],
             l=7.0,t=1.42,cw=[4.0,1.8,2.05],rh=0.27,fs=8.5)
    for r in [4,5]:
        for c in range(3): set_cell_bg(sh2.table.cell(r,c), LPURP_H)
    footer(sl,'Patient-level mean pooling ensures no data leakage and a stable per-patient representation')

    # ── Slide 3: Results ──────────────────────────────────────────────────────
    sl = prs.slides.add_slide(BLANK)
    rect(sl,0,0,13.33,7.5,'F9FBFD')
    header(sl,'Fusion Results — All Strategies',
           'Random Forest  ·  5-fold patient-level CV  ·  997 matched patients  ·  CD=0 UC=1')

    txt(sl,'Complete Comparison',0.25,1.08,13,0.30,size=11,bold=True,color=PURP)
    tbl_rows = [['Strategy','Dim','AUC','AP','Accuracy','CD F1','UC F1']]
    cat_colors = {'imaging': LIGHT_H, 'rna': LTEAL_H, 'fusion': LPURP_H,
                  'site': AMBER_H, 'ablation': GREY_H}
    row_cats = []
    for display, csv_name, src, dim, cat in STRATEGIES:
        df = dfs[src]
        r = strat(df, csv_name)
        if len(r) == 0: continue
        a, as_, p, ps, c, cs = summ(r)
        tbl_rows.append([display, dim,
                         f'{a:.3f}±{as_:.3f}', f'{p:.3f}±{ps:.3f}',
                         f'{c*100:.1f}%', f'{r["cd_f1"].mean():.3f}',
                         f'{r["uc_f1"].mean():.3f}'])
        row_cats.append(cat)
    sh = table(sl, tbl_rows, l=0.25, t=1.42,
               cw=[3.6,1.8,2.4,2.3,2.1,1.5,1.5], rh=0.28, fs=8.5)
    for ri, cat in enumerate(row_cats):
        bg = cat_colors.get(cat, GREY_H)
        for ci in range(7): set_cell_bg(sh.table.cell(ri+1, ci), bg)
    # bold the two best fusion rows
    for ri in [3,4]:
        for ci in range(7):
            for p in sh.table.cell(ri+1,ci).text_frame.paragraphs:
                for r in p.runs: r.font.bold = True

    divider(sl,5.10)
    txt(sl,'AUC Summary by Category',0.25,5.18,13,0.28,size=10,bold=True,color=PURP)
    bar_items = [
        ('concat raw/scaled', raw_s[0],   raw_s[1],   PURP_H),
        ('RNA only',          rna_s[0],   rna_s[1],   TEAL_H),
        ('late fusion',       s('late_fusion_base_rna')[0], s('late_fusion_base_rna')[1], '9B59B6'),
        ('concat pca128',     pca_c_s[0], pca_c_s[1], '8E44AD'),
        ('imaging only',      img_s[0],   img_s[1],   MID_H),
        ('site-matched',      s('site_matched_scaled_base')[0], s('site_matched_scaled_base')[1], ORANGE_H),
    ]
    max_auc=0.94; bscale=5.8; bx0=1.2
    for i,(name,auc,sd,col) in enumerate(bar_items):
        ly = 5.50 + i*0.30
        txt(sl,name,0.25,ly+0.04,0.9,0.22,size=7.5,color=DGREY)
        bw = auc/max_auc*bscale
        rect(sl,bx0,ly+0.02,bw,0.22,col)
        txt(sl,f'{auc:.3f}±{sd:.3f}',bx0+bw+0.06,ly+0.04,1.3,0.22,size=8.5,bold=True,color=DARK)
    rect(sl,bx0,7.12,bscale,0.02,MGREY_H)
    for tick in [0.7,0.8,0.9]:
        bxi = bx0+tick/max_auc*bscale
        txt(sl,str(tick),bxi-0.15,7.14,0.4,0.18,size=7,color=DGREY,align=PP_ALIGN.CENTER)

    footer(sl,'Purple = feature fusion  ·  Teal = RNA unimodal  ·  Blue = imaging unimodal  ·  Orange = site-matched')

    # ── Slide 4: Ablation ─────────────────────────────────────────────────────
    sl = prs.slides.add_slide(BLANK)
    rect(sl,0,0,13.33,7.5,'F9FBFD')
    header(sl,'Ablation Analysis','PCA compression effect · raw vs scaled concatenation')

    # left: PCA variance table
    txt(sl,'PCA-128 Unimodal Ablation',0.25,1.08,6.3,0.30,size=11,bold=True,color=PURP)
    sh=table(sl,[['Modality','Full dim','AUC (full)','PCA-128 AUC','ΔAUC','Var. retained'],
                 ['Imaging (prism2_base)','2,560',
                  f'{img_s[0]:.3f}',f'{img_pca_s[0]:.3f}',
                  f'{img_pca_s[0]-img_s[0]:+.3f}',f'{pca_var_img*100:.1f}%'],
                 ['RNA VST','17,963',
                  f'{rna_s[0]:.3f}',f'{rna_pca_s[0]:.3f}',
                  f'{rna_pca_s[0]-rna_s[0]:+.3f}',f'{pca_var_rna*100:.1f}%'],
                 ['Concat (img+RNA)','20,523',
                  f'{raw_s[0]:.3f}',f'{pca_c_s[0]:.3f}',
                  f'{pca_c_s[0]-raw_s[0]:+.3f}','—']],
             l=0.25,t=1.42,cw=[2.7,1.4,1.8,1.8,1.3,1.5],rh=0.30,fs=9)
    for c in range(6): set_cell_bg(sh.table.cell(2,c), LPURP_H)
    for p in sh.table.cell(2,c).text_frame.paragraphs:
        for r in p.runs: r.font.bold=True

    # visual diagnosis
    for i,(label,full_v,pca_v,col_f,col_p) in enumerate([
        ('Imaging: 99.2% retained → small loss', img_s[0],  img_pca_s[0], MID_H,  LIGHT_H),
        ('RNA: 81.0% retained → larger loss',    rna_s[0],  rna_pca_s[0], TEAL_H, LTEAL_H),
        ('Concat (img+RNA): driven by RNA loss',  raw_s[0],  pca_c_s[0],  PURP_H, LPURP_H),
    ]):
        ly = 2.50 + i*0.72
        rect(sl,0.25,ly,6.3,0.64,GREY_H,MGREY_H,0.3)
        txt(sl,label,0.35,ly+0.03,6.1,0.26,size=8.5,bold=True,color=DARK)
        scale=5.5; bx=0.35; max_v=0.96
        bw_f = full_v/max_v*scale
        bw_p = pca_v/max_v*scale
        rect(sl,bx,ly+0.32,bw_f,0.18,col_f)
        rect(sl,bx,ly+0.34,bw_p,0.14,col_p)
        txt(sl,f'full: {full_v:.3f}',bx+bw_f+0.05,ly+0.30,1.0,0.16,size=7.5,color=DARK)
        txt(sl,f'pca:  {pca_v:.3f}',bx+bw_p+0.05,ly+0.44,1.0,0.16,size=7.5,color=DGREY)

    divider(sl,4.72)
    txt(sl,'Raw vs Scaled Concatenation',0.25,4.82,6.3,0.28,size=11,bold=True,color=PURP)
    rect(sl,0.25,5.16,6.3,1.55,LPURP_H,PURP_H)
    mtxt(sl,['concat_raw == concat_scaled  (identical fold-for-fold)',
             '',
             'Random Forest uses rank-based splits (Gini impurity). StandardScaler',
             'preserves the ordering of values within a feature — so RF makes',
             'the exact same splits with or without scaling.',
             '',
             'RNA\'s 17,963 features do NOT dominate imaging\'s 2,560 features via scale.',
             'They dominate only because they are more predictive.'],
         0.38,5.22,6.0,1.42,size=8.5,color=DARK)

    # right: per-fold comparison
    txt(sl,'Per-Fold Detail: concat_raw vs concat_scaled vs concat_pca128',
        6.85,1.08,6.3,0.30,size=10,bold=True,color=PURP)
    raw_folds = abl[abl['strategy']=='concat_raw']
    cs_folds  = abl[abl['strategy']=='concat_scaled']
    pca_folds = mm[mm['strategy']=='concat_pca128_base_rna']
    table(sl,[['Fold','raw AUC','scaled AUC','pca128 AUC','Δraw-pca']]+
             [[str(f),
               f"{raw_folds[raw_folds['fold']==f]['auc'].values[0]:.3f}",
               f"{cs_folds[cs_folds['fold']==f]['auc'].values[0]:.3f}",
               f"{pca_folds[pca_folds['fold']==f]['auc'].values[0]:.3f}",
               f"{raw_folds[raw_folds['fold']==f]['auc'].values[0]-pca_folds[pca_folds['fold']==f]['auc'].values[0]:+.3f}"]
              for f in range(5)],
          l=6.85,t=1.42,cw=[1.0,1.8,1.8,1.8,1.8],rh=0.28,fs=9)

    txt(sl,'Per-Fold Detail: img_base_pca128 vs rna_pca128',
        6.85,3.55,6.3,0.28,size=10,bold=True,color=PURP)
    img_pca_folds = abl[abl['strategy']=='img_base_pca128']
    rna_pca_folds = abl[abl['strategy']=='rna_pca128']
    img_full_folds = mm[mm['strategy']=='img_base_patmean']
    rna_full_folds = mm[mm['strategy']=='rna_patmean']
    table(sl,[['Fold','img full','img pca128','rna full','rna pca128']]+
             [[str(f),
               f"{img_full_folds[img_full_folds['fold']==f]['auc'].values[0]:.3f}",
               f"{img_pca_folds[img_pca_folds['fold']==f]['auc'].values[0]:.3f}",
               f"{rna_full_folds[rna_full_folds['fold']==f]['auc'].values[0]:.3f}",
               f"{rna_pca_folds[rna_pca_folds['fold']==f]['auc'].values[0]:.3f}"]
              for f in range(5)],
          l=6.85,t=3.90,cw=[0.9,1.85,1.85,1.85,1.85],rh=0.28,fs=9)

    footer(sl,'PCA-128 retains 99.2% of imaging variance but only 81.0% of RNA variance  ·  RNA loss explains the concat_pca128 drop')

    # ── Slide 5: Discussion ───────────────────────────────────────────────────
    sl = prs.slides.add_slide(BLANK)
    rect(sl,0,0,13.33,7.5,'F9FBFD')
    header(sl,'Discussion & Next Steps',
           'What fusion adds (and does not add) for CD vs UC prediction')

    txt(sl,'Key Findings',0.25,1.08,6.3,0.30,size=11,bold=True,color=PURP)
    findings=[
        ('Fusion does not beat RNA alone',
         f'Best fusion (concat_raw) AUC {raw_s[0]:.3f} ≈ RNA alone {rna_s[0]:.3f}. '
         f'ΔAUC = {raw_s[0]-rna_s[0]:+.4f}. RNA already captures nearly all discriminative signal.'),
        ('RF is scale-invariant',
         'concat_raw == concat_scaled fold-for-fold. No need to scale features for Random Forest. '
         'Dimensionality imbalance only matters if features differ in predictive power, not scale.'),
        ('PCA hurts RNA more than imaging',
         f'PCA-128 retains {pca_var_img*100:.1f}% of imaging variance (ΔAUC {img_pca_s[0]-img_s[0]:+.3f}) '
         f'but only {pca_var_rna*100:.1f}% of RNA variance (ΔAUC {rna_pca_s[0]-rna_s[0]:+.3f}). '
         'RNA compression drives the concat_pca128 performance drop.'),
        ('Patient mean > site-matched',
         f'Site-matched fusion AUC {s("site_matched_scaled_base")[0]:.3f} vs patient-mean {raw_s[0]:.3f}. '
         'CD/UC is a patient-level label — averaging all sites reduces noise and improves the feature.'),
    ]
    y = 1.42
    for title,body in findings:
        rect(sl,0.25,y,6.3,0.72,LPURP_H,PURP_H)
        txt(sl,title,0.38,y+0.03,6.1,0.25,size=9,bold=True,color=PURP)
        txt(sl,body,0.38,y+0.28,6.1,0.40,size=8.5,color=DARK)
        y += 0.78

    txt(sl,'Limitations & Caveats',6.85,1.08,6.3,0.30,size=11,bold=True,color=ORANGE)
    caveats = [
        ('RF cannot model cross-modal interactions',
         'Concatenation lets RF select features from both modalities but cannot learn '
         'multiplicative or conditional interactions between RNA and imaging.'),
        ('Biopsy-site confound unresolved',
         'At-20-cm overrepresented (838/993 patients). Site-specific expression '
         'patterns may inflate RNA AUC beyond pure CD/UC signal.'),
        ('Patient-mean discards spatial information',
         'Mean pooling over slides loses the within-patient spatial distribution. '
         'Attention pooling (MIL) may recover site-specific signal.'),
        ('No probability calibration',
         'RF scores are not calibrated probabilities. Averaging uncalibrated scores '
         'in late fusion may not give the optimal combination.'),
    ]
    y = 1.42
    for title,body in caveats:
        rect(sl,6.85,y,6.3,0.72,AMBER_H,ORANGE_H,0.75)
        txt(sl,title,6.98,y+0.03,6.1,0.25,size=9,bold=True,color=ORANGE)
        txt(sl,body,6.98,y+0.28,6.1,0.40,size=8.5,color=DARK)
        y += 0.78

    divider(sl,5.36)
    txt(sl,'Recommended Next Steps',0.25,5.46,13,0.30,size=11,bold=True,color=PURP)
    for i,s_text in enumerate([
        '1. Neural late fusion — learned weighting of RNA + imaging scores',
        '2. Cross-modal attention — attend imaging patches conditioned on RNA embedding',
        '3. Site-covariate correction — remove at-20-cm bias before fusion',
        '4. SHAP importance — which genes vs. imaging dims drive the model',
        '5. Calibration — Platt scaling on RF probabilities before fusion',
    ]):
        lx = 0.25 + (i%3)*4.36; ly = 5.80 + (i//3)*0.52
        rect(sl,lx,ly,4.2,0.44,LPURP_H,PURP_H)
        txt(sl,s_text,lx+0.1,ly+0.06,4.0,0.34,size=8,color=DARK)

    footer(sl,'IBD Plexus / SPARC  ·  Virchow2 TRIDENT 0.3.0  ·  CombatSeq VST GSF1491805  ·  August 11, 2026')

    out_path, _ = next_versioned_path(os.path.join(REPORTS_DIR, 'multimodal_slides.pptx'))
    prs.save(out_path)
    print(f'PPTX saved: {out_path}')
    return out_path


if __name__ == '__main__':
    parser = argparse.ArgumentParser()
    parser.add_argument('--desc',
                        default='Multimodal fusion study: 6 strategies + PCA/scale ablations on 997-patient matched cohort',
                        help='One-line description for VERSIONS.md')
    args = parser.parse_args()

    pdf_path  = build_pdf()
    pptx_path = build_pptx()
    log_version(pdf_path,  args.desc)
    log_version(pptx_path, args.desc)
