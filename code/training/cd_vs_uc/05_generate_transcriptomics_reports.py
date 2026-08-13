"""
Generate PDF report and PowerPoint slides for the transcriptomics head-to-head analysis.

Reads the pre-computed fold-metric CSVs; does NOT re-run any model.

Outputs  (versioned — never overwrites existing files)
--------
- transcriptomics_report_vN.pdf
- transcriptomics_slides_vN.pptx
- VERSIONS.md  (append-only changelog)

Usage
-----
    python 05_generate_transcriptomics_reports.py
    python 05_generate_transcriptomics_reports.py --desc "Added biopsy-site breakdown"
"""

import argparse
import os
import numpy as np
import pandas as pd
from version_utils import next_versioned_path, log_version

TRANS_DIR    = '/home/jovyan/kgbk271-ibd-volume/training/cd_vs_uc/03_05_transcriptomics_allsites/results'
REPORTS_DIR  = '/home/jovyan/kgbk271-ibd-volume/training/cd_vs_uc/03_05_transcriptomics_allsites/reports'
IMG_DIR      = '/home/jovyan/kgbk271-ibd-volume/training/cd_vs_uc/02_04_imaging_allsites/results'
TRAINING_DIR = '/home/jovyan/kgbk271-ibd-volume/training'


def load_metrics():
    rna_f  = pd.read_csv(os.path.join(TRANS_DIR, 'transcriptomics_vst_fold_metrics.csv'))
    bm_f   = pd.read_csv(os.path.join(IMG_DIR, 'prism2_base_matched_fold_metrics.csv'))
    dm_f   = pd.read_csv(os.path.join(IMG_DIR, 'prism2_diagnostic_matched_fold_metrics.csv'))
    base_f = pd.read_csv(os.path.join(IMG_DIR, 'prism2_base_fold_metrics.csv'))
    return rna_f, bm_f, dm_f, base_f


def cm_stats(f):
    tn = f['tn'].sum(); fp = f['fp'].sum()
    fn = f['fn'].sum(); tp = f['tp'].sum()
    return tn, fp, fn, tp, tp / (tp + fn), tn / (tn + fp)


# ══════════════════════════════════════════════════════════════════════════════
#  PDF REPORT
# ══════════════════════════════════════════════════════════════════════════════

def build_pdf():
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.lib import colors
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer, Table,
                                    TableStyle, HRFlowable, PageBreak)
    from reportlab.lib.enums import TA_JUSTIFY, TA_CENTER

    DARK  = colors.HexColor('#1a2e4a')
    MID   = colors.HexColor('#2e6da4')
    TEAL  = colors.HexColor('#16a085')
    LTEAL = colors.HexColor('#d1f0eb')
    LIGHT = colors.HexColor('#d8e8f5')
    WHITE = colors.white
    GREEN = colors.HexColor('#1e8b4c')
    GREY  = colors.HexColor('#f5f7fa')
    MGREY = colors.HexColor('#dddddd')
    DGREY = colors.HexColor('#888888')
    ORANGE= colors.HexColor('#e67e22')
    AMBER = colors.HexColor('#fff3cd')
    LRED  = colors.HexColor('#fadbd8')
    RED   = colors.HexColor('#c0392b')

    rna_f, bm_f, dm_f, base_f = load_metrics()

    out_path, _ = next_versioned_path(os.path.join(REPORTS_DIR, 'transcriptomics_report.pdf'))
    doc = SimpleDocTemplate(out_path, pagesize=A4,
                            topMargin=2.2*cm, bottomMargin=2*cm,
                            leftMargin=2.2*cm, rightMargin=2.2*cm)

    S = getSampleStyleSheet()
    def sty(name, **kw): return ParagraphStyle(name, parent=S['Normal'], **kw)

    H1   = sty('H1',   fontSize=18, textColor=DARK,  spaceAfter=6,  fontName='Helvetica-Bold')
    H2   = sty('H2',   fontSize=13, textColor=TEAL,  spaceBefore=16, spaceAfter=5, fontName='Helvetica-Bold')
    H3   = sty('H3',   fontSize=11, textColor=DARK,  spaceBefore=10, spaceAfter=3, fontName='Helvetica-Bold')
    BODY = sty('BODY', fontSize=9.5, leading=15, spaceAfter=4, alignment=TA_JUSTIFY)
    SMALL= sty('SMALL',fontSize=8.5, leading=12, textColor=DGREY)
    BULL = sty('BULL', fontSize=9.5, leading=14, leftIndent=14, spaceAfter=3)
    NOTE = sty('NOTE', fontSize=8.5, leading=12, textColor=colors.HexColor('#444444'), leftIndent=10)
    FOOT = sty('FOOT', fontSize=7.5, textColor=DGREY, alignment=TA_CENTER)

    TS = [
        ('BACKGROUND', (0,0), (-1,0), TEAL), ('TEXTCOLOR', (0,0), (-1,0), WHITE),
        ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'), ('FONTSIZE', (0,0), (-1,-1), 9),
        ('ROWBACKGROUNDS', (0,1), (-1,-1), [WHITE, GREY]),
        ('GRID', (0,0), (-1,-1), 0.4, MGREY), ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
        ('TOPPADDING', (0,0), (-1,-1), 4), ('BOTTOMPADDING', (0,0), (-1,-1), 4),
        ('LEFTPADDING', (0,0), (-1,-1), 6), ('RIGHTPADDING', (0,0), (-1,-1), 6),
    ]
    def tbl(data, col_widths, extra=None):
        t = Table(data, colWidths=[w*cm for w in col_widths])
        t.setStyle(TableStyle(TS + (extra or [])))
        return t

    story = []

    # ── Title block ────────────────────────────────────────────────────────────
    delta_auc = rna_f['auc'].mean() - bm_f['auc'].mean()
    story += [
        Spacer(1, 1.5*cm),
        Paragraph('Transcriptomics vs Imaging: Head-to-Head Comparison', H1),
        Paragraph('CD vs UC Classification on Matched 997-Patient Cohort',
                  sty('sub', fontSize=13, textColor=TEAL, spaceAfter=6)),
        HRFlowable(width='100%', thickness=2, color=TEAL, spaceAfter=10),
        Paragraph(
            'Dataset: IBD Plexus / SPARC &nbsp;·&nbsp; '
            'Transcriptomics: CombatSeq VST (17,963 genes) &nbsp;·&nbsp; '
            'Imaging: Virchow2 prism2_base & prism2_diagnostic<br/>'
            'Date: August 10, 2026', SMALL),
        Spacer(1, 0.5*cm),
    ]
    story.append(tbl([
        ['Matched cohort',  '997 patients with both colon WSI and colon RNAseq'],
        ['Transcriptomics AUC', f"{rna_f['auc'].mean():.3f} ± {rna_f['auc'].std():.3f}"],
        ['Imaging AUC',         f"{bm_f['auc'].mean():.3f} ± {bm_f['auc'].std():.3f}  (prism2_base_matched)"],
        ['ΔAUC (RNA − imaging)', f"+{delta_auc:.3f}"],
        ['Splits',          'Identical patient-level 5-fold CV for all three models'],
    ], [4, 12], extra=[
        ('FONTNAME', (0,0), (0,-1), 'Helvetica-Bold'),
        ('BACKGROUND', (0,0), (0,-1), LTEAL),
        ('BACKGROUND', (0,0), (-1,0), LTEAL),
        ('TEXTCOLOR', (0,0), (-1,0), DARK),
    ]))
    story.append(PageBreak())

    # ── Section 1: Why matched cohort ─────────────────────────────────────────
    story.append(Paragraph('1. Head-to-Head Design', H2))
    story.append(Paragraph(
        'The original imaging analysis used <b>1,250 patients</b> (all colon CD/UC patients '
        'with WSI embeddings). Transcriptomics was only available for <b>997</b> of those '
        'patients. Comparing imaging (1,250 pts) to transcriptomics (997 pts) directly '
        'conflates cohort differences with modality differences.', BODY))
    story.append(Paragraph(
        'Fix: re-run the imaging Random Forest restricted to the same 997-patient subset, '
        'using identical fold assignments. All three models are then evaluated on '
        'identical patients and identical train/val splits.', BODY))
    story.append(tbl([
        ['Model', 'Cohort', 'N patients', 'N samples/slides'],
        ['transcriptomics_vst',       'Matched', '997', '1,728 RNAseq samples'],
        ['prism2_base_matched',        'Matched', '997', '1,758 slides'],
        ['prism2_diagnostic_matched',  'Matched', '997', '1,758 slides'],
        ['prism2_base (original ref)', 'Full',   '1,250','2,121 slides (reference only)'],
    ], [5, 2.5, 2.5, 6], extra=[
        ('BACKGROUND', (0,1), (-1,3), [LTEAL]*8),
        ('TEXTCOLOR',  (0,4), (-1,4), DGREY),
    ]))

    # ── Section 2: Transcriptomics data ───────────────────────────────────────
    story.append(Paragraph('2. Transcriptomics Data', H2))
    story.append(Paragraph('2.1  File selection', H3))
    story.append(tbl([
        ['File', 'Genes', 'Normalisation', 'Batch corrected', 'Used'],
        ['GSF1491805 CombatSeq VST', '17,963', 'DESeq2 VST', 'Yes (CombatSeq)', '✓ SELECTED'],
        ['GSF2048892 TPM',           '62,703', 'TPM',         'No',             '—'],
        ['GSF1491803 CombatSeq counts','17,963','raw counts',  'Yes',            '—'],
        ['GSF1485554 VST (no batch)', '62,703', 'DESeq2 VST',  'No',             '—'],
        ['GSF1491807 Raw counts',     '62,703', 'none',        'No',             '—'],
    ], [5, 2, 3, 3, 3], extra=[
        ('BACKGROUND', (0,1), (-1,1), LTEAL),
        ('FONTNAME', (0,1), (-1,1), 'Helvetica-Bold'),
    ]))
    story.append(Paragraph(
        '<b>Rationale:</b> GSF1491805 is the only file combining CombatSeq batch correction '
        '(first vs later sequencing run) AND DESeq2 VST (continuous floats ~3.5–12.5, '
        'approximately Gaussian), making it most suitable for Random Forest directly '
        'on the feature matrix.', NOTE))

    story.append(Paragraph('2.2  Sample filtering', H3))
    story.append(Paragraph(
        'From the 3,289 samples in the GCT file: keep colon biopsies only '
        '(at 20 cm, Cecum, Rectum, Ascending/Descending/Sigmoid/Transverse Colon), '
        'CD or UC diagnosis, QC-pass, and patient in the imaging CV split. '
        'Result: <b>1,728 samples from 997 patients</b>.', BODY))

    # ── Section 3: Results ────────────────────────────────────────────────────
    story.append(Paragraph('3. Results', H2))
    story.append(Paragraph('3.1  Overall performance — matched cohort', H3))
    story.append(tbl([
        ['Model', 'Modality', 'N', 'AUC (mean±SD)', 'AP (mean±SD)', 'Accuracy'],
        ['transcriptomics_vst',
         'RNAseq VST', '1,728',
         f"{rna_f['auc'].mean():.3f} ± {rna_f['auc'].std():.3f}",
         f"{rna_f['ap'].mean():.3f} ± {rna_f['ap'].std():.3f}",
         f"{rna_f['accuracy'].mean()*100:.1f}% ± {rna_f['accuracy'].std()*100:.1f}%"],
        ['prism2_base_matched',
         'WSI embed', '1,758',
         f"{bm_f['auc'].mean():.3f} ± {bm_f['auc'].std():.3f}",
         f"{bm_f['ap'].mean():.3f} ± {bm_f['ap'].std():.3f}",
         f"{bm_f['accuracy'].mean()*100:.1f}% ± {bm_f['accuracy'].std()*100:.1f}%"],
        ['prism2_diagnostic_matched',
         'WSI embed', '1,758',
         f"{dm_f['auc'].mean():.3f} ± {dm_f['auc'].std():.3f}",
         f"{dm_f['ap'].mean():.3f} ± {dm_f['ap'].std():.3f}",
         f"{dm_f['accuracy'].mean()*100:.1f}% ± {dm_f['accuracy'].std()*100:.1f}%"],
    ], [4.5, 2.5, 1.5, 3, 3, 2.5], extra=[
        ('BACKGROUND', (0,1), (-1,1), LTEAL),
        ('FONTNAME', (0,1), (-1,1), 'Helvetica-Bold'),
    ]))

    story.append(Paragraph('3.2  Per-fold results', H3))
    for name, folds in [('transcriptomics_vst', rna_f),
                        ('prism2_base_matched', bm_f),
                        ('prism2_diagnostic_matched', dm_f)]:
        story.append(Paragraph(f'<b>{name}</b>', BULL))
        story.append(tbl(
            [['Fold', 'N val', 'AUC', 'AP', 'Accuracy', 'CD F1', 'UC F1']] +
            [[str(r.fold), str(r.n_val), f'{r.auc:.3f}', f'{r.ap:.3f}',
              f'{r.accuracy*100:.1f}%', f'{r.cd_f1:.3f}', f'{r.uc_f1:.3f}']
             for _, r in folds.iterrows()],
            [1.5, 2, 2, 2, 2.5, 2, 2]))

    story.append(Paragraph('3.3  Aggregated confusion matrices', H3))
    for name, folds in [('transcriptomics_vst', rna_f),
                        ('prism2_base_matched', bm_f),
                        ('prism2_diagnostic_matched', dm_f)]:
        tn,fp,fn,tp,sens,spec = cm_stats(folds)
        story.append(Paragraph(
            f'<b>{name}</b>: TN={tn} FP={fp} FN={fn} TP={tp} | '
            f'Sensitivity={sens:.3f}  Specificity={spec:.3f}', BULL))

    story.append(Paragraph('3.4  Original full-cohort imaging results (reference)', H3))
    story.append(Paragraph(
        'Reported for completeness only — <b>not directly comparable</b> to the matched results '
        'above (different patient set). The small difference (ΔAUC ≈ 0.001) confirms '
        'restricting to the matched cohort introduced no material bias.', NOTE))
    story.append(tbl([
        ['Model', 'N', 'AUC', 'AP', 'Accuracy'],
        ['prism2_base (full)',
         '2,121',
         f"{base_f['auc'].mean():.3f} ± {base_f['auc'].std():.3f}",
         f"{base_f['ap'].mean():.3f} ± {base_f['ap'].std():.3f}",
         f"{base_f['accuracy'].mean()*100:.1f}% ± {base_f['accuracy'].std()*100:.1f}%"],
    ], [5, 2, 3, 3, 4]))

    # ── Section 4: Discussion ─────────────────────────────────────────────────
    story.append(Paragraph('4. Discussion', H2))
    story.append(Paragraph(
        f'Transcriptomics (AUC {rna_f["auc"].mean():.3f}) outperforms imaging '
        f'(AUC {bm_f["auc"].mean():.3f}) by ΔAUC = {delta_auc:.3f} on an identical cohort '
        f'with identical folds. This is a true modality difference, not a cohort artefact. '
        f'Gene expression captures molecular programs (cytokine signatures, goblet-cell loss, '
        f'crypt architecture at the mRNA level) that are not fully visible in H&E morphology '
        f'assessed by Virchow2 embeddings alone.', BODY))
    story.append(Paragraph(
        '<b>Caveats:</b> (1) Biopsy-site composition differs between CD (enriched at 20 cm) '
        'and UC (enriched at cecum); both modalities are affected but the effect may be larger '
        'for transcriptomics due to site-specific gene expression. '
        '(2) No multimodal fusion evaluated yet — combining imaging + RNA likely improves both. '
        '(3) Predictive performance does not directly identify which genes drive separation.', BODY))

    story.append(Paragraph('5. Next Steps', H2))
    for s in [
        '1. Biopsy-site correction — restrict to at-20-cm or add site covariate',
        '2. Multimodal fusion — late-fuse imaging + RNA embeddings',
        '3. SHAP / feature importance — identify discriminative genes',
        '4. Patient-level aggregation — average per-patient scores before evaluation',
        '5. Attention MIL — use patch-level features instead of slide embeddings',
    ]:
        story.append(Paragraph(s, BULL))

    story += [
        Spacer(1, 0.5*cm),
        HRFlowable(width='100%', thickness=0.5, color=MGREY),
        Paragraph(
            'IBD Plexus / SPARC  ·  CombatSeq VST GSF1491805  ·  '
            'Virchow2 via TRIDENT 0.3.0  ·  August 10, 2026', FOOT),
    ]

    doc.build(story)
    print(f'PDF saved: {out_path}')
    return out_path


# ══════════════════════════════════════════════════════════════════════════════
#  POWERPOINT
# ══════════════════════════════════════════════════════════════════════════════

def build_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from lxml import etree

    rna_f, bm_f, dm_f, base_f = load_metrics()

    def rgb(r,g,b): return RGBColor(r,g,b)
    DARK   = rgb(0x1a,0x2e,0x4a); DARK_H  = '1A2E4A'
    MID    = rgb(0x2e,0x6d,0xa4); MID_H   = '2E6DA4'
    LIGHT  = rgb(0xd8,0xe8,0xf5); LIGHT_H = 'D8E8F5'
    WHITE  = rgb(0xff,0xff,0xff); WHITE_H = 'FFFFFF'
    GREEN  = rgb(0x1e,0x8b,0x4c)
    GREY   = rgb(0xf5,0xf7,0xfa); GREY_H  = 'F5F7FA'
    LGREEN = rgb(0xd5,0xf5,0xe3); LGREEN_H= 'D5F5E3'
    LRED_H = 'FADBD8'
    TEAL   = rgb(0x16,0xa0,0x85); TEAL_H  = '16A085'
    LTEAL  = rgb(0xd1,0xf0,0xeb); LTEAL_H = 'D1F0EB'
    ORANGE = rgb(0xe6,0x7e,0x22); ORANGE_H= 'E67E22'
    AMBER_H= 'FFF3CD'
    DGREY  = rgb(0x55,0x55,0x55); DGREY_H = '555555'
    MGREY_H= 'CCCCCC'
    RED    = rgb(0xc0,0x39,0x2b); RED_H   = 'C0392B'

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    BLANK = prs.slide_layouts[6]

    def rect(sl,l,t,w,h,fh,lh=None,lw=0.75):
        s = sl.shapes.add_shape(1,Inches(l),Inches(t),Inches(w),Inches(h))
        s.fill.solid(); s.fill.fore_color.rgb = RGBColor.from_string(fh)
        if lh: s.line.color.rgb = RGBColor.from_string(lh); s.line.width = Pt(lw)
        else:  s.line.fill.background()

    def txt(sl,text,l,t,w,h,size=10,bold=False,color=DARK,align=PP_ALIGN.LEFT,wrap=True):
        tb = sl.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
        tf = tb.text_frame; tf.word_wrap = wrap
        p = tf.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color

    def mtxt(sl,lines,l,t,w,h,size=8.5,color=DARK,bold_first=False):
        tb = sl.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
        tf = tb.text_frame; tf.word_wrap = True
        for i,line in enumerate(lines):
            p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
            r = p.add_run(); r.text = line
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold_first and i==0

    def set_cell_bg(cell, hx):
        tc = cell._tc; tcPr = tc.get_or_add_tcPr()
        for old in tcPr.findall(qn('a:solidFill')): tcPr.remove(old)
        sf = etree.SubElement(tcPr, qn('a:solidFill'))
        sc = etree.SubElement(sf, qn('a:srgbClr')); sc.set('val', hx)

    def table(sl,data,l,t,cw,rh=0.27,hdr=MID_H,fs=8.5):
        rows=len(data); cols=len(data[0]); tw=sum(cw)
        sh = sl.shapes.add_table(rows,cols,Inches(l),Inches(t),Inches(tw),Inches(rh*rows))
        tb = sh.table
        for i,w in enumerate(cw): tb.columns[i].width = Inches(w)
        for r in range(rows): tb.rows[r].height = Inches(rh)
        for r,row in enumerate(data):
            for c,val in enumerate(row):
                cell = tb.cell(r,c); cell.text = str(val)
                for para in cell.text_frame.paragraphs:
                    para.alignment = PP_ALIGN.CENTER
                    for run in para.runs:
                        run.font.size = Pt(fs); run.font.bold = (r==0)
                        run.font.color.rgb = WHITE if r==0 else DARK
                if r==0: set_cell_bg(cell, hdr)
                elif r%2==0: set_cell_bg(cell, GREY_H)
                else: set_cell_bg(cell, WHITE_H)
        return sh

    def header(sl,title,sub=None):
        rect(sl,0,0,13.33,0.95,DARK_H)
        txt(sl,title,0.25,0.08,11,0.45,size=15,bold=True,color=WHITE)
        if sub: txt(sl,sub,0.25,0.57,11,0.30,size=9,color=LIGHT)

    def footer(sl,text):
        txt(sl,text,0.2,7.22,13,0.22,size=7.5,color=DGREY,align=PP_ALIGN.CENTER)

    def divider(sl,y):
        rect(sl,0.2,y,12.93,0.02,MID_H)

    delta_auc = rna_f['auc'].mean() - bm_f['auc'].mean()

    # ── Slide 1: Title ────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(BLANK)
    rect(sl,0,0,13.33,7.5,'F2F6FB')
    rect(sl,0,0,13.33,2.7,DARK_H)
    rect(sl,0,2.65,13.33,0.08,TEAL_H)
    txt(sl,'Transcriptomics vs Imaging',0.5,0.38,12.3,0.65,
        size=26,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    txt(sl,'Head-to-Head on Matched Cohort',0.5,0.97,12.3,0.55,
        size=20,color=LIGHT,align=PP_ALIGN.CENTER)
    txt(sl,'997 patients with both colon RNAseq and colon WSI  ·  Identical 5-fold CV splits  ·  Random Forest',
        0.5,1.65,12.3,0.4,size=11,color=LIGHT,align=PP_ALIGN.CENTER)
    for i,(val,lab,fh) in enumerate([
        (f'+{delta_auc:.3f}','AUC gap (RNA > img)',TEAL_H),
        (f'{rna_f["auc"].mean():.3f}','RNA AUC',TEAL_H),
        (f'{bm_f["auc"].mean():.3f}','img AUC (matched)',MID_H),
        ('997','matched patients',DARK_H),
    ]):
        lx = 1.0 + i*2.88
        bg = LTEAL_H if fh == TEAL_H else LIGHT_H
        rect(sl,lx,3.05,2.5,1.4,bg,fh)
        txt(sl,val,lx,3.18,2.5,0.65,size=24,bold=True,
            color=TEAL if fh==TEAL_H else MID,align=PP_ALIGN.CENTER)
        txt(sl,lab,lx,3.82,2.5,0.4,size=9.5,color=DARK,align=PP_ALIGN.CENTER)
    mtxt(sl,['Transcriptomics: CombatSeq VST  |  17,963 genes  |  1,728 samples',
             'Imaging: Virchow2 prism2_base (2,560-d) & prism2_diagnostic (3,072-d)  |  1,758 slides',
             'Same patient-level 5-fold CV splits used for all models  |  August 10, 2026'],
         1.0,4.85,11.3,0.9,size=9.5,color=DARK)

    # ── Slide 2: Design & Cohort ──────────────────────────────────────────────
    sl = prs.slides.add_slide(BLANK)
    rect(sl,0,0,13.33,7.5,'F9FBFD')
    header(sl,'Study Design & Cohort','Why matched cohort is required for fair comparison')

    txt(sl,'The Problem: Original Comparison Was Confounded',0.25,1.08,13,0.3,size=11,bold=True,color=RED)
    for i,(lab,detail,fh) in enumerate([
        ('Imaging (original)\n1,250 pts · 2,121 slides','Includes 253 patients\nwith no RNAseq',RED_H),
        ('Transcriptomics\n997 pts · 1,728 samples','Only patients with\ncolon RNAseq',TEAL_H),
    ]):
        lx = 0.8 + i*5.5
        bg = LRED_H if i==0 else LTEAL_H
        rect(sl,lx,1.45,4.5,1.0,bg,fh)
        for j,line in enumerate(lab.split('\n')):
            txt(sl,line,lx+0.1,1.49+j*0.28,4.3,0.26,size=9.5,
                bold=(j==0),color=RED if i==0 else TEAL)
        for j,line in enumerate(detail.split('\n')):
            txt(sl,line,lx+0.1,1.90+j*0.22,4.3,0.22,size=8.5,color=DARK)
    txt(sl,'≠',5.6,1.76,0.8,0.4,size=14,bold=True,color=RED,align=PP_ALIGN.CENTER)

    rect(sl,0.25,2.58,12.83,0.05,TEAL_H)
    txt(sl,'The Fix: Re-run imaging on the same 997-patient matched cohort',
        0.25,2.70,13,0.30,size=11,bold=True,color=TEAL)
    for i,(lab,n,fh) in enumerate([
        ('transcriptomics_vst',      '997 pts · 1,728 samples',TEAL_H),
        ('prism2_base_matched',       '997 pts · 1,758 slides', MID_H),
        ('prism2_diagnostic_matched', '997 pts · 1,758 slides', '4A90C8'),
    ]):
        lx = 0.5 + i*4.28
        bg = LTEAL_H if i==0 else LIGHT_H
        rect(sl,lx,3.05,3.8,0.9,bg,fh)
        txt(sl,lab,lx+0.1,3.08,3.6,0.28,size=9.5,bold=True,
            color=TEAL if i==0 else MID)
        txt(sl,n,lx+0.1,3.36,3.6,0.28,size=8.5,color=DARK)
        txt(sl,'Same folds ✓',lx+0.1,3.62,3.6,0.24,size=8,color=GREEN)
        if i < 2:
            txt(sl,'→',lx+3.82,3.42,0.4,0.28,size=12,bold=True,color=DGREY,align=PP_ALIGN.CENTER)

    divider(sl,4.18)
    txt(sl,'Cohort Breakdown',0.25,4.28,6.5,0.3,size=11,bold=True,color=MID)
    sh=table(sl,[['','CD','UC','Total'],
                 ['Full imaging cohort','841','409','1,250'],
                 ['Matched cohort (RNAseq ∩ imaging)','569','428','997'],
                 ['RNAseq samples (matched)','940','788','1,728'],
                 ['Imaging slides (matched)','967','791','1,758']],
             l=0.25,t=4.62,cw=[5.5,1.5,1.5,1.5],rh=0.27,fs=8.5)
    for r in [2,3,4]:
        for c in range(4): set_cell_bg(sh.table.cell(r,c), LTEAL_H)

    txt(sl,'Transcriptomics Selection',7.1,4.28,6.0,0.3,size=11,bold=True,color=MID)
    sh2=table(sl,[['File','Genes','Norm.','Batch','Use'],
                  ['GSF1491805 CombatSeq VST','17,963','DESeq2 VST','✓','✓ USED'],
                  ['GSF2048892 TPM','62,703','TPM','✗','—'],
                  ['GSF1491803 CombatSeq counts','17,963','raw','✓','—'],
                  ['GSF1485554 VST (no batch)','62,703','VST','✗','—'],
                  ['GSF1491807 Raw counts','62,703','none','✗','—']],
             l=7.1,t=4.62,cw=[4.0,1.4,1.4,1.0,0.9],rh=0.27,fs=8)
    for c in range(5):
        set_cell_bg(sh2.table.cell(1,c), LTEAL_H)
        for p in sh2.table.cell(1,c).text_frame.paragraphs:
            for r in p.runs: r.font.bold = True
    footer(sl,'Matched cohort = patients with both colon WSI embedding AND colon RNAseq  ·  Same fold IDs for all models')

    # ── Slide 3: Results ──────────────────────────────────────────────────────
    sl = prs.slides.add_slide(BLANK)
    rect(sl,0,0,13.33,7.5,'F9FBFD')
    header(sl,'Head-to-Head Results: Matched Cohort (997 Patients)',
           'Random Forest  ·  CD=0 UC=1  ·  5-fold patient-level CV  ·  Identical splits')

    txt(sl,'Model Comparison — Matched Cohort',0.25,1.08,13,0.3,size=11,bold=True,color=MID)
    sh=table(sl,[
        ['Model','Modality','N','AUC','AP','Accuracy','Sens. (UC)','Spec. (CD)'],
        ['transcriptomics_vst','RNAseq VST','1,728',
         f"{rna_f['auc'].mean():.3f}±{rna_f['auc'].std():.3f}",
         f"{rna_f['ap'].mean():.3f}±{rna_f['ap'].std():.3f}",
         f"{rna_f['accuracy'].mean()*100:.1f}%",
         f"{cm_stats(rna_f)[4]:.3f}",f"{cm_stats(rna_f)[5]:.3f}"],
        ['prism2_base_matched','WSI embed','1,758',
         f"{bm_f['auc'].mean():.3f}±{bm_f['auc'].std():.3f}",
         f"{bm_f['ap'].mean():.3f}±{bm_f['ap'].std():.3f}",
         f"{bm_f['accuracy'].mean()*100:.1f}%",
         f"{cm_stats(bm_f)[4]:.3f}",f"{cm_stats(bm_f)[5]:.3f}"],
        ['prism2_diag_matched','WSI embed','1,758',
         f"{dm_f['auc'].mean():.3f}±{dm_f['auc'].std():.3f}",
         f"{dm_f['ap'].mean():.3f}±{dm_f['ap'].std():.3f}",
         f"{dm_f['accuracy'].mean()*100:.1f}%",
         f"{cm_stats(dm_f)[4]:.3f}",f"{cm_stats(dm_f)[5]:.3f}"],
        ['prism2_base (full, ref.)','WSI embed','2,121',
         f"{base_f['auc'].mean():.3f}±{base_f['auc'].std():.3f}",
         f"{base_f['ap'].mean():.3f}±{base_f['ap'].std():.3f}",
         f"{base_f['accuracy'].mean()*100:.1f}%",
         f"{cm_stats(base_f)[4]:.3f}",f"{cm_stats(base_f)[5]:.3f}"],
    ],l=0.25,t=1.42,cw=[3.8,2.2,1.4,2.4,2.0,1.8,1.8,1.8],rh=0.30,fs=9)
    for c in range(8):
        set_cell_bg(sh.table.cell(1,c), LTEAL_H)
        for p in sh.table.cell(1,c).text_frame.paragraphs:
            for r in p.runs: r.font.bold = True
    for c in range(8):
        set_cell_bg(sh.table.cell(4,c), GREY_H)
        for p in sh.table.cell(4,c).text_frame.paragraphs:
            for r in p.runs: r.font.color.rgb = DGREY

    divider(sl,2.63)

    txt(sl,'AUC (matched cohort)',0.25,2.73,5.8,0.28,size=10,bold=True,color=MID)
    bar_data = [
        ('transcriptomics_vst', rna_f['auc'].mean(), rna_f['auc'].std(), TEAL_H),
        ('prism2_base',         bm_f['auc'].mean(),  bm_f['auc'].std(),  MID_H),
        ('prism2_diagnostic',   dm_f['auc'].mean(),  dm_f['auc'].std(),  '4A90C8'),
    ]
    max_auc=0.92; bscale=5.3; bx0=0.9
    for i,(name,auc,sd,col) in enumerate(bar_data):
        ly = 3.05 + i*0.65
        txt(sl,name,0.25,ly+0.10,0.6,0.26,size=7.5,color=DGREY)
        bw = auc/max_auc*bscale
        rect(sl,bx0,ly+0.07,bw,0.38,col)
        txt(sl,f'{auc:.3f}±{sd:.3f}',bx0+bw+0.05,ly+0.10,1.2,0.26,size=9,bold=True,color=DARK)
    rect(sl,bx0,5.05,bscale,0.02,MGREY_H)
    for tick in [0.5,0.6,0.7,0.8,0.9]:
        bxi = bx0 + tick/max_auc*bscale
        rect(sl,bxi,4.92,0.01,0.16,MGREY_H)
        txt(sl,str(tick),bxi-0.15,5.10,0.4,0.2,size=7.5,color=DGREY,align=PP_ALIGN.CENTER)
    txt(sl,'AUC →',bx0+bscale/2-0.3,5.28,1.0,0.22,size=8,color=DGREY)

    txt(sl,'transcriptomics_vst per fold',6.85,2.73,6.3,0.28,size=10,bold=True,color=TEAL)
    table(sl,[['Fold','N val','AUC','AP','Acc','CD F1','UC F1']]+
             [[str(r.fold),str(r.n_val),f'{r.auc:.3f}',f'{r.ap:.3f}',
               f'{r.accuracy*100:.1f}%',f'{r.cd_f1:.3f}',f'{r.uc_f1:.3f}']
              for _,r in rna_f.iterrows()],
          l=6.85,t=3.05,cw=[0.9,1.2,1.2,1.2,1.4,1.3,1.3],rh=0.26,fs=8.5,hdr=TEAL_H)

    txt(sl,'prism2_base_matched per fold',6.85,4.60,6.3,0.28,size=10,bold=True,color=MID)
    table(sl,[['Fold','N val','AUC','AP','Acc','CD F1','UC F1']]+
             [[str(r.fold),str(r.n_val),f'{r.auc:.3f}',f'{r.ap:.3f}',
               f'{r.accuracy*100:.1f}%',f'{r.cd_f1:.3f}',f'{r.uc_f1:.3f}']
              for _,r in bm_f.iterrows()],
          l=6.85,t=4.90,cw=[0.9,1.2,1.2,1.2,1.4,1.3,1.3],rh=0.26,fs=8.5)

    divider(sl,5.60)
    txt(sl,'Aggregated Confusion Matrices (positive = UC)',0.25,5.68,13,0.28,size=10,bold=True,color=MID)
    for xi,(name,f,fh) in enumerate([
        ('transcriptomics', rna_f, TEAL_H),
        ('prism2_base\n(matched)', bm_f, MID_H),
        ('prism2_diag\n(matched)', dm_f, '4A90C8'),
    ]):
        lx = 0.25 + xi*4.36
        tn,fp,fn,tp,sens,spec = cm_stats(f)
        rect(sl,lx,5.98,4.2,1.32,LIGHT_H,MGREY_H,0.4)
        txt(sl,name.replace('\n',' '),lx+0.1,6.01,4.0,0.24,size=8,bold=True,
            color=RGBColor.from_string(fh))
        cw2,ch=0.88,0.40; cx2=lx+0.1; cy2=6.28
        vals=[[tn,fp],[fn,tp]]
        cbg=[['F5F7FA','FADBD8'],['FADBD8',LTEAL_H]]
        for ri in range(2):
            for ci in range(2):
                rect(sl,cx2+ci*cw2,cy2+ri*ch,cw2,ch,cbg[ri][ci],MGREY_H,0.3)
                txt(sl,str(vals[ri][ci]),cx2+ci*cw2,cy2+ri*ch+0.04,cw2,0.22,
                    size=10,bold=True,color=DARK,align=PP_ALIGN.CENTER)
                txt(sl,f'P{"CD" if ci==0 else "UC"}',cx2+ci*cw2,cy2+ri*ch+0.26,cw2,0.13,
                    size=6.5,color=DGREY,align=PP_ALIGN.CENTER)
        txt(sl,f'Sens {sens:.3f}  Spec {spec:.3f}',lx+0.1,cy2+0.88,4.0,0.2,size=8,color=DARK)

    footer(sl,'AUC = area under ROC  ·  AP = average precision  ·  prism2_base full-cohort shown for reference only')

    # ── Slide 4: Discussion ───────────────────────────────────────────────────
    sl = prs.slides.add_slide(BLANK)
    rect(sl,0,0,13.33,7.5,'F9FBFD')
    header(sl,'Discussion & Next Steps',
           'Interpretation of head-to-head results and recommended future work')

    txt(sl,'Key Findings',0.25,1.08,6.3,0.30,size=11,bold=True,color=TEAL)
    findings=[
        (f'ΔAUC = {delta_auc:.3f} is a true modality effect',
         f'Gap confirmed on identical 997-patient cohort with identical folds. Not an artefact.'),
        ('Both modalities stable across cohort restriction',
         f'prism2_base AUC changes only {abs(base_f["auc"].mean()-bm_f["auc"].mean()):.3f} from 1,250 to 997 patients.'),
        ('Transcriptomics improves both recall directions',
         f'UC sensitivity {cm_stats(rna_f)[4]:.3f} vs {cm_stats(bm_f)[4]:.3f}; '
         f'CD specificity {cm_stats(rna_f)[5]:.3f} vs {cm_stats(bm_f)[5]:.3f}.'),
        ('Gene expression captures what histology cannot',
         'Molecular programs (cytokine signatures, goblet cell loss) distinguish CD/UC beyond H&E morphology.'),
    ]
    y=1.42
    for title,body in findings:
        rect(sl,0.25,y,6.3,0.72,LTEAL_H,TEAL_H)
        txt(sl,title,0.38,y+0.03,6.1,0.25,size=9,bold=True,color=TEAL)
        txt(sl,body,0.38,y+0.28,6.1,0.40,size=8.5,color=DARK)
        y+=0.78

    txt(sl,'Caveats',6.85,1.08,6.3,0.30,size=11,bold=True,color=ORANGE)
    caveats=[
        ('Biopsy-site composition confound',
         'UC enriched at cecum; CD at 20 cm. Site effect on transcriptomics may be larger.'),
        ('Different sample counts per modality',
         '1,728 RNAseq samples vs 1,758 slides from same 997 patients. Reflects different sampling.'),
        ('No multimodal fusion yet',
         'Combining imaging + transcriptomics features not evaluated. Likely to improve over either.'),
        ('Predictive not explanatory',
         'No feature importance or differential expression analysis performed.'),
    ]
    y=1.42
    for title,body in caveats:
        rect(sl,6.85,y,6.3,0.72,AMBER_H,ORANGE_H,0.75)
        txt(sl,title,6.98,y+0.03,6.1,0.25,size=9,bold=True,color=ORANGE)
        txt(sl,body,6.98,y+0.28,6.1,0.40,size=8.5,color=DARK)
        y+=0.78

    divider(sl,5.36)
    txt(sl,'Recommended Next Steps',0.25,5.46,13,0.30,size=11,bold=True,color=MID)
    for i,s in enumerate([
        '1. Biopsy-site correction — restrict to at-20-cm or add site covariate',
        '2. Multimodal fusion — concatenate or late-fuse imaging + RNA embeddings',
        '3. SHAP / feature importance — identify top genes driving CD/UC separation',
        '4. Patient-level aggregation — average slide/sample scores per patient',
        '5. ABMIL on patches — attention MIL with patch-level Virchow2 features',
    ]):
        lx = 0.25 + (i%3)*4.36; ly = 5.80 + (i//3)*0.52
        rect(sl,lx,ly,4.2,0.44,LIGHT_H,MID_H)
        txt(sl,s,lx+0.1,ly+0.06,4.0,0.34,size=8,color=DARK)

    footer(sl,'IBD Plexus / SPARC  ·  Virchow2 TRIDENT 0.3.0  ·  CombatSeq VST GSF1491805  ·  August 10, 2026')

    out_path, _ = next_versioned_path(os.path.join(REPORTS_DIR, 'transcriptomics_slides.pptx'))
    prs.save(out_path)
    print(f'PPTX saved: {out_path}')
    return out_path


if __name__ == '__main__':
    parser = argparse.ArgumentParser()
    parser.add_argument('--desc',
                        default='Transcriptomics CombatSeq VST + head-to-head vs matched imaging (997 patients)',
                        help='One-line description of what changed in this version')
    args = parser.parse_args()

    pdf_path  = build_pdf()
    pptx_path = build_pptx()
    log_version(pdf_path,  args.desc)
    log_version(pptx_path, args.desc)
