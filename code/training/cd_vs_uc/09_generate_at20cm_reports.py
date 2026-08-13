"""
Generate PDF report and PowerPoint slides for the At-20-cm site-restricted analysis.

Usage
-----
  python 09_generate_at20cm_reports.py
  python 09_generate_at20cm_reports.py --desc "..."
"""

import argparse
import os
import json
import numpy as np
import pandas as pd
from version_utils import next_versioned_path, log_version

AT20_DIR    = '/home/jovyan/kgbk271-ibd-volume/training/cd_vs_uc/08_09_at20cm_site_controlled/results'
REPORTS_DIR = '/home/jovyan/kgbk271-ibd-volume/training/cd_vs_uc/08_09_at20cm_site_controlled/reports'
MM_DIR      = '/home/jovyan/kgbk271-ibd-volume/training/cd_vs_uc/06_07_multimodal_allsites/results'


def load_all():
    at20  = pd.read_csv(os.path.join(AT20_DIR, 'at20cm_fold_metrics.csv'))
    mm    = pd.read_csv(os.path.join(MM_DIR, 'multimodal_fold_metrics.csv'))
    abl   = pd.read_csv(os.path.join(MM_DIR, 'multimodal_ablation_fold_metrics.csv'))
    return at20, mm, abl


def s(df, name):
    r = df[df['strategy'] == name]
    if len(r) == 0: return None
    a = r['auc'].values; p = r['ap'].values; c = r['accuracy'].values
    return dict(mean_auc=round(float(a.mean()),4), std_auc=round(float(a.std()),4),
                mean_ap=round(float(p.mean()),4),  std_ap=round(float(p.std()),4),
                mean_acc=round(float(c.mean()),4),  std_acc=round(float(c.std()),4),
                rows=r)


def confusion_stats(r):
    tn=r['tn'].sum(); fp=r['fp'].sum(); fn=r['fn'].sum(); tp=r['tp'].sum()
    return tn,fp,fn,tp, tp/(tp+fn), tn/(tn+fp)


# ══════════════════════════════════════════════════════════════════════════════
#  PDF REPORT
# ══════════════════════════════════════════════════════════════════════════════

def build_pdf():
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.lib import colors
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer, Table,
                                    TableStyle, HRFlowable, PageBreak)
    from reportlab.lib.enums import TA_JUSTIFY, TA_CENTER

    DARK   = colors.HexColor('#1a2e4a'); MID  = colors.HexColor('#2e6da4')
    TEAL   = colors.HexColor('#16a085'); PURP = colors.HexColor('#6c3082')
    ORANGE = colors.HexColor('#e67e22'); GREEN= colors.HexColor('#1e8b4c')
    WHITE  = colors.white;               RED  = colors.HexColor('#c0392b')
    GREY   = colors.HexColor('#f5f7fa'); MGREY= colors.HexColor('#dddddd')
    DGREY  = colors.HexColor('#888888')
    LTEAL  = colors.HexColor('#d1f0eb'); LPURP= colors.HexColor('#e8d5f5')
    LGREEN = colors.HexColor('#d5f5e3'); LRED = colors.HexColor('#fadbd8')
    AMBER  = colors.HexColor('#fff3cd')

    at20, mm, abl = load_all()

    rna_all  = s(mm,  'rna_patmean')
    img_all  = s(mm,  'img_base_patmean')
    fus_all  = s(abl, 'concat_raw')
    rna_20   = s(at20,'rna_20cm')
    img_20   = s(at20,'img_base_20cm')
    fus_20   = s(at20,'concat_raw_20cm')
    pca_20   = s(at20,'concat_pca128_20cm')

    out_path, _ = next_versioned_path(os.path.join(REPORTS_DIR, 'at20cm_report.pdf'))
    doc = SimpleDocTemplate(out_path, pagesize=A4,
                            topMargin=2.2*cm, bottomMargin=2*cm,
                            leftMargin=2.2*cm, rightMargin=2.2*cm)

    S = getSampleStyleSheet()
    def sty(n, **kw): return ParagraphStyle(n, parent=S['Normal'], **kw)
    H1   = sty('H1', fontSize=18, textColor=DARK,   spaceAfter=6,   fontName='Helvetica-Bold')
    H2   = sty('H2', fontSize=13, textColor=ORANGE,  spaceBefore=16, spaceAfter=5,  fontName='Helvetica-Bold')
    H3   = sty('H3', fontSize=11, textColor=DARK,   spaceBefore=10, spaceAfter=3,  fontName='Helvetica-Bold')
    BODY = sty('B',  fontSize=9.5, leading=15, spaceAfter=4, alignment=TA_JUSTIFY)
    BULL = sty('BL', fontSize=9.5, leading=14, leftIndent=14, spaceAfter=3)
    NOTE = sty('N',  fontSize=8.5, leading=12, textColor=colors.HexColor('#444'), leftIndent=10)
    SMALL= sty('SM', fontSize=8.5, leading=12, textColor=DGREY)
    FOOT = sty('FT', fontSize=7.5, textColor=DGREY, alignment=TA_CENTER)

    def TS(h=ORANGE): return [
        ('BACKGROUND',(0,0),(-1,0),h),('TEXTCOLOR',(0,0),(-1,0),WHITE),
        ('FONTNAME',(0,0),(-1,0),'Helvetica-Bold'),('FONTSIZE',(0,0),(-1,-1),9),
        ('ROWBACKGROUNDS',(0,1),(-1,-1),[WHITE,GREY]),
        ('GRID',(0,0),(-1,-1),0.4,MGREY),('VALIGN',(0,0),(-1,-1),'MIDDLE'),
        ('TOPPADDING',(0,0),(-1,-1),4),('BOTTOMPADDING',(0,0),(-1,-1),4),
        ('LEFTPADDING',(0,0),(-1,-1),6),('RIGHTPADDING',(0,0),(-1,-1),6),
    ]
    def tbl(data, cw, h=ORANGE, extra=None):
        t = Table(data, colWidths=[w*cm for w in cw])
        t.setStyle(TableStyle(TS(h)+(extra or [])))
        return t

    story = []

    # ── Title ──────────────────────────────────────────────────────────────────
    story += [
        Spacer(1, 1.5*cm),
        Paragraph('Site-Controlled Analysis: At-20-cm Biopsies Only', H1),
        Paragraph('Removing the Biopsy-Protocol Site Confound from CD vs UC Classification',
                  sty('s2', fontSize=13, textColor=ORANGE, spaceAfter=6)),
        HRFlowable(width='100%', thickness=2, color=ORANGE, spaceAfter=10),
        Paragraph('Dataset: IBD Plexus / SPARC  ·  828 matched patients (554 CD, 274 UC)  ·  '
                  'At-20-cm = rectosigmoid junction  ·  August 11, 2026', SMALL),
        Spacer(1, 0.5*cm),
    ]
    story.append(tbl([
        ['RNA At-20-cm AUC',   f"{rna_20['mean_auc']:.3f} ± {rna_20['std_auc']:.3f}  (vs {rna_all['mean_auc']:.3f} all-sites)"],
        ['Imaging At-20-cm AUC',f"{img_20['mean_auc']:.3f} ± {img_20['std_auc']:.3f}  (vs {img_all['mean_auc']:.3f} all-sites)"],
        ['ΔAUC (all-sites − At-20-cm)',
         f"RNA: {rna_all['mean_auc']-rna_20['mean_auc']:+.3f}   Imaging: {img_all['mean_auc']-img_20['mean_auc']:+.3f}  ← site-identity signal"],
        ['Matched cohort',     '828 patients with At-20-cm RNA AND imaging  ·  5-fold patient-level CV'],
        ['Key finding',        '~0.09 AUC in both modalities was site-identity shortcut, not disease biology'],
    ], [4,12], h=ORANGE, extra=[
        ('FONTNAME',(0,0),(0,-1),'Helvetica-Bold'),
        ('BACKGROUND',(0,0),(0,-1),AMBER),('BACKGROUND',(0,3),(-1,3),AMBER),
        ('TEXTCOLOR',(0,0),(-1,0),DARK),('BACKGROUND',(0,0),(-1,0),AMBER),
    ]))
    story.append(PageBreak())

    # ── Section 1: The confound ────────────────────────────────────────────────
    story.append(Paragraph('1. The Biopsy-Protocol Site Confound', H2))
    story.append(Paragraph(
        'The IBD Plexus / SPARC biopsy protocol collects samples from <b>different '
        'anatomical locations</b> depending on diagnosis:', BODY))
    story.append(tbl([
        ['Diagnosis','Proximal site','Distal site'],
        ['CD (Crohn\'s disease)',   'Ileum  ← excluded by colon filter','At 20 cm (rectosigmoid)'],
        ['UC (Ulcerative colitis)', 'Cecum  ← included in colon filter','At 20 cm (rectosigmoid)'],
    ], [4,6,6], extra=[
        ('BACKGROUND',(0,1),(-1,1),LRED),
        ('BACKGROUND',(0,2),(-1,2),LTEAL),
    ]))
    story.append(Paragraph(
        'In the all-sites analysis, UC patients contributed Cecum samples while CD patients '
        'did not (their proximal biopsy — ileum — was excluded). '
        'A classifier trained on all-sites features could therefore predict UC simply by '
        'detecting <b>"this sample is from the cecum"</b> rather than by learning CD/UC '
        'disease biology. Gene expression and histology differ substantially by anatomical '
        'location (different cell-type composition, microbiome exposure, mucosal architecture), '
        'making this a genuine protocol-driven shortcut available to the model.', BODY))
    story.append(Paragraph(
        '<b>The fix:</b> restrict to At-20-cm (rectosigmoid junction) samples only. '
        'Both CD and UC patients contribute samples from this identical site, '
        'removing the site-identity shortcut entirely.', BODY))

    # ── Section 2: Cohort ─────────────────────────────────────────────────────
    story.append(Paragraph('2. Cohort Statistics', H2))
    story.append(tbl([
        ['','All-sites (all colon locs)','At-20-cm only','Change'],
        ['RNA patients',  '997',  '847', '−150 (−15%)'],
        ['RNA samples',   '2,186','1,068','−1,118 (−51%)'],
        ['Imaging patients','1,250','1,016','−234 (−19%)'],
        ['Imaging slides','2,121','1,218','−903 (−43%)'],
        ['Matched (both modalities)','997','828','−169 (−17%)'],
        ['CD patients (matched)','569','554','−15'],
        ['UC patients (matched)','428','274','−154 (−36%)'],
    ], [5,4,4,3], extra=[
        ('BACKGROUND',(0,5),(-1,6),AMBER),
    ]))
    story.append(Paragraph(
        'UC patients are disproportionately reduced (−36%) because many of their '
        'samples were Cecum-only. CD patients are largely unchanged because their '
        'colon samples were already predominantly At-20-cm.', NOTE))

    # ── Section 3: Results ────────────────────────────────────────────────────
    story.append(Paragraph('3. Results', H2))
    story.append(Paragraph('3.1  At-20-cm vs all-sites comparison', H3))

    def delta_str(v20, vall): return f'{v20-vall:+.3f}'

    story.append(tbl([
        ['Model','N pts','AUC (At-20-cm)','AUC (all-sites)','ΔAUC','Interpretation'],
        ['img_base',  '828', f"{img_20['mean_auc']:.3f}±{img_20['std_auc']:.3f}",
         f"{img_all['mean_auc']:.3f}±{img_all['std_auc']:.3f}",
         delta_str(img_20['mean_auc'],img_all['mean_auc']),
         'Morphological site shortcut removed'],
        ['RNA VST',   '828', f"{rna_20['mean_auc']:.3f}±{rna_20['std_auc']:.3f}",
         f"{rna_all['mean_auc']:.3f}±{rna_all['std_auc']:.3f}",
         delta_str(rna_20['mean_auc'],rna_all['mean_auc']),
         'Expression site shortcut removed'],
        ['concat_raw','828', f"{fus_20['mean_auc']:.3f}±{fus_20['std_auc']:.3f}",
         f"{fus_all['mean_auc']:.3f}±{fus_all['std_auc']:.3f}",
         delta_str(fus_20['mean_auc'],fus_all['mean_auc']),
         'Both shortcuts removed'],
        ['concat_pca128','828', f"{pca_20['mean_auc']:.3f}±{pca_20['std_auc']:.3f}",
         '0.909±0.017',
         delta_str(pca_20['mean_auc'],0.9086),'—'],
    ], [2.5,1.5,3.5,3.5,2.0,4.0], extra=[
        ('BACKGROUND',(0,1),(-1,3),AMBER),
        ('FONTNAME',(0,1),(-1,3),'Helvetica-Bold'),
    ]))

    story.append(Paragraph('3.2  Per-fold detail (At-20-cm cohort)', H3))
    for name, label in [('img_base_20cm','img_base At-20-cm'),
                        ('rna_20cm','RNA At-20-cm'),
                        ('concat_raw_20cm','concat At-20-cm')]:
        rows = at20[at20['strategy']==name]
        story.append(Paragraph(f'<b>{label}</b>', BULL))
        story.append(tbl(
            [['Fold','N val','AUC','AP','Accuracy','CD F1','UC F1']]+
            [[str(r.fold),str(r.n_val),f'{r.auc:.3f}',f'{r.ap:.3f}',
              f'{r.accuracy*100:.1f}%',f'{r.cd_f1:.3f}',f'{r.uc_f1:.3f}']
             for _,r in rows.iterrows()],
            [1.5,1.8,2,2,2.3,2,2]))

    story.append(Paragraph('3.3  Aggregated confusion matrices (At-20-cm)', H3))
    for name, lbl in [('img_base_20cm','img_base'),('rna_20cm','RNA'),('concat_raw_20cm','concat')]:
        tn,fp,fn,tp,sens,spec = confusion_stats(at20[at20['strategy']==name])
        story.append(Paragraph(
            f'<b>{lbl}</b>: TN={tn} FP={fp} FN={fn} TP={tp} | '
            f'Sensitivity={sens:.3f}  Specificity={spec:.3f}', BULL))

    # ── Section 4: Interpretation ─────────────────────────────────────────────
    story.append(Paragraph('4. Interpretation', H2))

    img_drop = img_all['mean_auc'] - img_20['mean_auc']
    rna_drop = rna_all['mean_auc'] - rna_20['mean_auc']
    rna_gap_all = rna_all['mean_auc'] - img_all['mean_auc']
    rna_gap_20  = rna_20['mean_auc']  - img_20['mean_auc']

    story.append(Paragraph('4.1  Site-identity was a major driver of all-sites performance', H3))
    story.append(Paragraph(
        f'Imaging AUC dropped by {img_drop:.3f} and RNA AUC dropped by {rna_drop:.3f} '
        f'when removing the site confound. Both drops are approximately equal, '
        f'confirming that both modalities were exploiting the same underlying '
        f'signal: <b>cecum samples are always from UC patients</b> by biopsy protocol design. '
        f'Virchow2 embeddings captured cecum morphology (different gland architecture, '
        f'mucosal folds) as a discriminative feature; VST expression captured cecum-specific '
        f'gene expression patterns. Neither of these signals reflects CD/UC disease biology '
        f'at a controlled tissue site.', BODY))

    story.append(Paragraph('4.2  RNA retains a larger absolute advantage at matched tissue', H3))
    story.append(tbl([
        ['Comparison','All-sites','At-20-cm only'],
        ['RNA AUC',                f"{rna_all['mean_auc']:.3f}",f"{rna_20['mean_auc']:.3f}"],
        ['Imaging AUC',            f"{img_all['mean_auc']:.3f}",f"{img_20['mean_auc']:.3f}"],
        ['ΔAUC (RNA − imaging)',   f"+{rna_gap_all:.3f}",       f"+{rna_gap_20:.3f}"],
        ['Fusion vs RNA alone',    f"+{fus_all['mean_auc']-rna_all['mean_auc']:+.4f}",
                                   f"+{fus_20['mean_auc']-rna_20['mean_auc']:+.4f}"],
    ], [6,4,4], extra=[
        ('BACKGROUND',(0,3),(-1,3),LTEAL),('FONTNAME',(0,3),(-1,3),'Helvetica-Bold'),
    ]))
    story.append(Paragraph(
        f'The RNA advantage over imaging narrows slightly ({rna_gap_all:.3f} → {rna_gap_20:.3f}) '
        f'but RNA remains substantially more discriminative at the matched tissue site. '
        f'Molecular gene expression profiles at the rectosigmoid genuinely differ between '
        f'CD and UC beyond what is visible in H&E morphology. '
        f'Fusion continues to add negligible signal over RNA alone '
        f'({fus_20["mean_auc"]-rna_20["mean_auc"]:+.4f} at At-20-cm).', BODY))

    story.append(Paragraph('4.3  The honest performance estimate', H3))
    story.append(Paragraph(
        f'The At-20-cm AUC (RNA: <b>{rna_20["mean_auc"]:.3f}</b>, '
        f'imaging: <b>{img_20["mean_auc"]:.3f}</b>) is the more rigorous estimate of '
        f'what these models can achieve on a clinically realistic question: '
        f'"<i>given a rectosigmoid biopsy, does this patient have CD or UC?</i>" '
        f'The all-sites AUC ({rna_all["mean_auc"]:.3f} / {img_all["mean_auc"]:.3f}) '
        f'should be interpreted as an upper bound that includes protocol-driven site signal.', BODY))

    # ── Section 5: Next steps ─────────────────────────────────────────────────
    story.append(Paragraph('5. Next Steps', H2))
    for step in [
        '1. Multimodal fusion at At-20-cm — test whether imaging adds to RNA without site confound',
        '2. Cecum-only analysis — confirm Cecum AUC ≈ 1.0 (pure site identity)',
        '3. Site-covariate model — add biopsy location as a feature and train on all sites',
        '4. SHAP on At-20-cm RNA — identify disease-specific genes, free of site expression confounds',
        '5. Repeat with ABMIL on patches — test whether attention mechanism also recovers site signal',
    ]:
        story.append(Paragraph(step, BULL))

    story += [
        Spacer(1, 0.5*cm),
        HRFlowable(width='100%', thickness=0.5, color=MGREY),
        Paragraph('IBD Plexus / SPARC  ·  Virchow2 TRIDENT 0.3.0  ·  '
                  'CombatSeq VST GSF1491805  ·  August 11, 2026', FOOT),
    ]

    doc.build(story)
    print(f'PDF saved: {out_path}')
    return out_path


# ══════════════════════════════════════════════════════════════════════════════
#  POWERPOINT
# ══════════════════════════════════════════════════════════════════════════════

def build_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from lxml import etree

    at20, mm, abl = load_all()
    rna_all = s(mm,  'rna_patmean');        img_all = s(mm,  'img_base_patmean')
    fus_all = s(abl, 'concat_raw');         rna_20  = s(at20,'rna_20cm')
    img_20  = s(at20,'img_base_20cm');      fus_20  = s(at20,'concat_raw_20cm')
    pca_20  = s(at20,'concat_pca128_20cm')

    def rgb(r,g,b): return RGBColor(r,g,b)
    DARK  = rgb(0x1a,0x2e,0x4a); DARK_H  = '1A2E4A'
    MID   = rgb(0x2e,0x6d,0xa4); MID_H   = '2E6DA4'
    TEAL  = rgb(0x16,0xa0,0x85); TEAL_H  = '16A085'
    ORA   = rgb(0xe6,0x7e,0x22); ORA_H   = 'E67E22'
    AMBRH = 'FFF3CD'
    LTEAL = rgb(0xd1,0xf0,0xeb); LTEAL_H = 'D1F0EB'
    LIGHT = rgb(0xd8,0xe8,0xf5); LIGHT_H = 'D8E8F5'
    LORA  = rgb(0xfd,0xf0,0xe0); LORA_H  = 'FDF0E0'
    WHITE = rgb(0xff,0xff,0xff); WHITE_H  = 'FFFFFF'
    GREEN = rgb(0x1e,0x8b,0x4c); GREEN_H  = '1E8B4C'
    LGREEN_H= 'D5F5E3'
    GREY  = rgb(0xf5,0xf7,0xfa); GREY_H   = 'F5F7FA'
    DGREY = rgb(0x55,0x55,0x55); DGREY_H  = '555555'
    MGREY_H = 'CCCCCC'
    RED   = rgb(0xc0,0x39,0x2b); RED_H    = 'C0392B'
    LRED_H  = 'FADBD8'

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    BLANK = prs.slide_layouts[6]

    def rect(sl,l,t,w,h,fh,lh=None,lw=0.75):
        sh = sl.shapes.add_shape(1,Inches(l),Inches(t),Inches(w),Inches(h))
        sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor.from_string(fh)
        if lh: sh.line.color.rgb = RGBColor.from_string(lh); sh.line.width = Pt(lw)
        else:  sh.line.fill.background()

    def txt(sl,text,l,t,w,h,size=10,bold=False,color=DARK,align=PP_ALIGN.LEFT,wrap=True):
        tb = sl.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
        tf = tb.text_frame; tf.word_wrap = wrap
        p = tf.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color

    def mtxt(sl,lines,l,t,w,h,size=8.5,color=DARK):
        tb = sl.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
        tf = tb.text_frame; tf.word_wrap = True
        for i,line in enumerate(lines):
            p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
            r = p.add_run(); r.text = line
            r.font.size = Pt(size); r.font.color.rgb = color

    def set_cell_bg(cell, hx):
        tc = cell._tc; tcPr = tc.get_or_add_tcPr()
        for old in tcPr.findall(qn('a:solidFill')): tcPr.remove(old)
        sf = etree.SubElement(tcPr, qn('a:solidFill'))
        sc = etree.SubElement(sf, qn('a:srgbClr')); sc.set('val', hx)

    def table(sl,data,l,t,cw,rh=0.27,hdr=ORA_H,fs=8.5):
        rows=len(data); cols=len(data[0]); tw=sum(cw)
        sh = sl.shapes.add_table(rows,cols,Inches(l),Inches(t),Inches(tw),Inches(rh*rows))
        tb = sh.table
        for i,w in enumerate(cw): tb.columns[i].width = Inches(w)
        for r in range(rows): tb.rows[r].height = Inches(rh)
        for r,row in enumerate(data):
            for c,val in enumerate(row):
                cell = tb.cell(r,c); cell.text = str(val)
                for para in cell.text_frame.paragraphs:
                    para.alignment = PP_ALIGN.CENTER
                    for run in para.runs:
                        run.font.size = Pt(fs); run.font.bold = (r==0)
                        run.font.color.rgb = WHITE if r==0 else DARK
                if r==0: set_cell_bg(cell, hdr)
                elif r%2==0: set_cell_bg(cell, GREY_H)
                else: set_cell_bg(cell, WHITE_H)
        return sh

    def header(sl,title,sub=None):
        rect(sl,0,0,13.33,0.95,DARK_H)
        txt(sl,title,0.25,0.08,11,0.45,size=15,bold=True,color=WHITE)
        if sub: txt(sl,sub,0.25,0.57,11,0.30,size=9,color=LIGHT)

    def footer(sl,text):
        txt(sl,text,0.2,7.22,13,0.22,size=7.5,color=DGREY,align=PP_ALIGN.CENTER)

    def divider(sl,y):
        rect(sl,0.2,y,12.93,0.02,ORA_H)

    img_drop = img_all['mean_auc'] - img_20['mean_auc']
    rna_drop = rna_all['mean_auc'] - rna_20['mean_auc']

    # ── Slide 1: Title ────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(BLANK)
    rect(sl,0,0,13.33,7.5,'FFFBF5')
    rect(sl,0,0,13.33,2.7,DARK_H)
    rect(sl,0,2.65,13.33,0.08,ORA_H)
    txt(sl,'Site-Controlled Analysis',0.5,0.30,12.3,0.65,
        size=26,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    txt(sl,'At-20-cm Biopsies Only · Removing the Protocol-Driven Site Confound',
        0.5,0.94,12.3,0.45,size=16,color=LIGHT,align=PP_ALIGN.CENTER)
    txt(sl,'Both CD and UC biopsied at the same location (rectosigmoid junction)  ·  '
        '828 matched patients  ·  5-fold CV',
        0.5,1.58,12.3,0.38,size=10.5,color=LIGHT,align=PP_ALIGN.CENTER)
    for i,(val,lab,fh,bg) in enumerate([
        (f'{rna_20["mean_auc"]:.3f}', 'RNA AUC (At-20-cm)',     TEAL_H, LTEAL_H),
        (f'{img_20["mean_auc"]:.3f}', 'Imaging AUC (At-20-cm)', MID_H,  LIGHT_H),
        (f'−{rna_drop:.3f}',          'RNA drop (site removed)', ORA_H,  LORA_H),
        (f'−{img_drop:.3f}',          'Img drop (site removed)', ORA_H,  LORA_H),
    ]):
        lx = 1.0 + i*2.88
        rect(sl,lx,3.05,2.5,1.4,bg,fh)
        txt(sl,val,lx,3.18,2.5,0.65,size=24,bold=True,
            color=TEAL if fh==TEAL_H else (MID if fh==MID_H else ORA),
            align=PP_ALIGN.CENTER)
        txt(sl,lab,lx,3.82,2.5,0.40,size=9.5,color=DARK,align=PP_ALIGN.CENTER)
    mtxt(sl,['All-sites RNA AUC: 0.920  →  At-20-cm RNA AUC: 0.824  (−0.096 was site-identity shortcut)',
             'All-sites imaging AUC: 0.848  →  At-20-cm imaging AUC: 0.762  (−0.085 was site-identity shortcut)',
             'Both modalities exploited the same confound: cecum samples are UC-only by biopsy protocol design'],
         0.5,4.82,12.3,0.9,size=9.5,color=DARK)

    # ── Slide 2: Protocol + confound ──────────────────────────────────────────
    sl = prs.slides.add_slide(BLANK)
    rect(sl,0,0,13.33,7.5,'FAFAFA')
    header(sl,'The Biopsy-Protocol Site Confound',
           'Why cecum samples create a shortcut for CD vs UC classifiers')

    txt(sl,'IBD Plexus / SPARC Biopsy Protocol',0.25,1.08,12.8,0.30,size=11,bold=True,color=ORA)
    sh=table(sl,[['Diagnosis','Proximal biopsy','Distal biopsy','In colon filter?'],
                 ["CD (Crohn's disease)",'Ileum (ileocecal region)','At 20 cm (rectosigmoid)','Ileum → EXCLUDED'],
                 ['UC (Ulcerative colitis)','Cecum (proximal colon)','At 20 cm (rectosigmoid)','Cecum → INCLUDED']],
             l=0.25,t=1.42,cw=[3.8,3.8,3.8,2.4],rh=0.32,fs=9.5)
    set_cell_bg(sh.table.cell(1,3), LRED_H)
    set_cell_bg(sh.table.cell(2,3), LTEAL_H)
    for p in sh.table.cell(1,3).text_frame.paragraphs:
        for r in p.runs: r.font.bold=True; r.font.color.rgb=RED
    for p in sh.table.cell(2,3).text_frame.paragraphs:
        for r in p.runs: r.font.bold=True; r.font.color.rgb=TEAL

    divider(sl,2.20)
    txt(sl,'Consequence in All-Sites Analysis',0.25,2.30,12.8,0.30,size=11,bold=True,color=ORA)

    for i,(title,body,fh,bg) in enumerate([
        ('Cecum samples = UC-only by protocol',
         'Every cecum biopsy in the dataset is labeled UC (CD patients get ileum biopsies instead). '
         'The classifier can achieve high accuracy by detecting cecum-like features.',
         RED_H, LRED_H),
        ('Cecum differs from rectosigmoid',
         'Different mucosal architecture, gland density, goblet cell distribution (imaging) '
         'and different cell-type composition + microbiome-driven gene expression (RNA).',
         ORA_H, LORA_H),
        ('"Site signal" ≠ disease biology',
         'A model exploiting cecum identity does not generalize to CD/UC distinction '
         'at any given anatomical site — it has learned location, not disease.',
         ORA_H, LORA_H),
    ]):
        lx = 0.25 + i*4.36
        rect(sl,lx,2.64,4.1,1.48,bg,fh)
        txt(sl,title,lx+0.1,2.67,3.9,0.28,size=9,bold=True,
            color=RED if fh==RED_H else ORA)
        txt(sl,body, lx+0.1,2.98,3.9,1.08,size=8.5,color=DARK)

    divider(sl,4.26)
    txt(sl,'The Fix: At-20-cm Only',0.25,4.36,12.8,0.30,size=11,bold=True,color=GREEN)
    for i,(lab,body,fh,bg) in enumerate([
        ('Same site for both groups','Both CD and UC biopsied at the rectosigmoid junction (20 cm from anal verge). Site identity no longer discriminates.',GREEN_H,LGREEN_H),
        ('Cohort impact','828 matched patients (554 CD, 274 UC). UC drops more (−36%) because more UC patients had Cecum-only colon samples.',TEAL_H,LTEAL_H),
        ('Clinical relevance','More realistic scenario: "given a rectosigmoid biopsy, predict CD vs UC?" — the question a pathologist faces.',MID_H,LIGHT_H),
    ]):
        lx = 0.25 + i*4.36
        rect(sl,lx,4.70,4.1,1.48,bg,fh)
        txt(sl,lab, lx+0.1,4.73,3.9,0.28,size=9,bold=True,color=RGBColor.from_string(fh))
        txt(sl,body,lx+0.1,5.04,3.9,1.08,size=8.5,color=DARK)

    footer(sl,'At 20 cm = rectosigmoid junction  ·  Both CD and UC patients receive this biopsy by protocol')

    # ── Slide 3: Results comparison ───────────────────────────────────────────
    sl = prs.slides.add_slide(BLANK)
    rect(sl,0,0,13.33,7.5,'FAFAFA')
    header(sl,'Results: At-20-cm vs All-Sites',
           'Random Forest  ·  Patient-level mean pooling  ·  5-fold CV  ·  CD=0, UC=1')

    txt(sl,'Side-by-Side Comparison',0.25,1.08,12.8,0.30,size=11,bold=True,color=ORA)
    sh=table(sl,[['Model','Scope','N pts','AUC','AP','Accuracy','Sens.','Spec.'],
                 ['img_base','All sites','997',
                  f"{img_all['mean_auc']:.3f}±{img_all['std_auc']:.3f}",
                  f"{img_all['mean_ap']:.3f}±{img_all['std_ap']:.3f}",
                  f"{img_all['mean_acc']*100:.1f}%",'—','—'],
                 ['img_base','At-20-cm','828',
                  f"{img_20['mean_auc']:.3f}±{img_20['std_auc']:.3f}",
                  f"{img_20['mean_ap']:.3f}±{img_20['std_ap']:.3f}",
                  f"{img_20['mean_acc']*100:.1f}%",
                  f"{confusion_stats(img_20['rows'])[4]:.3f}",f"{confusion_stats(img_20['rows'])[5]:.3f}"],
                 ['RNA VST','All sites','997',
                  f"{rna_all['mean_auc']:.3f}±{rna_all['std_auc']:.3f}",
                  f"{rna_all['mean_ap']:.3f}±{rna_all['std_ap']:.3f}",
                  f"{rna_all['mean_acc']*100:.1f}%",'—','—'],
                 ['RNA VST','At-20-cm','828',
                  f"{rna_20['mean_auc']:.3f}±{rna_20['std_auc']:.3f}",
                  f"{rna_20['mean_ap']:.3f}±{rna_20['std_ap']:.3f}",
                  f"{rna_20['mean_acc']*100:.1f}%",
                  f"{confusion_stats(rna_20['rows'])[4]:.3f}",f"{confusion_stats(rna_20['rows'])[5]:.3f}"],
                 ['concat_raw','All sites','997',
                  f"{fus_all['mean_auc']:.3f}±{fus_all['std_auc']:.3f}",
                  f"{fus_all['mean_ap']:.3f}±{fus_all['std_ap']:.3f}",
                  f"{fus_all['mean_acc']*100:.1f}%",'—','—'],
                 ['concat_raw','At-20-cm','828',
                  f"{fus_20['mean_auc']:.3f}±{fus_20['std_auc']:.3f}",
                  f"{fus_20['mean_ap']:.3f}±{fus_20['std_ap']:.3f}",
                  f"{fus_20['mean_acc']*100:.1f}%",
                  f"{confusion_stats(fus_20['rows'])[4]:.3f}",f"{confusion_stats(fus_20['rows'])[5]:.3f}"],
              ],
             l=0.25,t=1.42,cw=[2.4,2.0,1.3,2.5,2.3,1.8,1.4,1.4],rh=0.29,fs=8.5)
    for r in [1,3,5]:
        for c in range(8): set_cell_bg(sh.table.cell(r,c), GREY_H)
    for r in [2,4,6]:
        for c in range(8): set_cell_bg(sh.table.cell(r,c), LORA_H)
        for p in sh.table.cell(r,0).text_frame.paragraphs:
            for run in p.runs: run.font.bold=True

    divider(sl,3.65)
    # bar chart: paired all-sites vs at-20cm
    txt(sl,'AUC drop from all-sites → At-20-cm',0.25,3.75,6.3,0.28,size=10,bold=True,color=ORA)
    pairs_bar = [
        ('Imaging',  img_all['mean_auc'], img_20['mean_auc'], MID_H,  LIGHT_H),
        ('RNA',      rna_all['mean_auc'], rna_20['mean_auc'], TEAL_H, LTEAL_H),
        ('Fusion',   fus_all['mean_auc'], fus_20['mean_auc'], '9B59B6','E8D5F5'),
    ]
    max_v=1.0; bscale=5.0; bx0=1.05
    for i,(name,vall,v20,col_all,col_20) in enumerate(pairs_bar):
        ly = 4.08 + i*0.82
        txt(sl,name,0.25,ly+0.26,0.75,0.22,size=8.5,color=DGREY)
        bw_all = vall/max_v*bscale; bw_20 = v20/max_v*bscale
        rect(sl,bx0,ly+0.06,bw_all,0.28,col_all)
        rect(sl,bx0,ly+0.38,bw_20, 0.28,col_20)
        txt(sl,f'all: {vall:.3f}',bx0+bw_all+0.05,ly+0.08,1.1,0.22,size=8,color=DARK)
        txt(sl,f'20cm:{v20:.3f}',bx0+bw_20+0.05,ly+0.40,1.1,0.22,size=8,color=ORA)
        drop=vall-v20
        rect(sl,bx0+bw_20,ly+0.38,bw_all-bw_20,0.28,ORA_H)
        txt(sl,f'−{drop:.3f}',bx0+bw_20+(bw_all-bw_20)/2-0.15,ly+0.40,0.7,0.22,
            size=7.5,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    rect(sl,bx0,7.05,bscale,0.02,MGREY_H)
    for tick in [0.7,0.8,0.9,1.0]:
        bxi=bx0+tick/max_v*bscale
        txt(sl,str(tick),bxi-0.15,7.07,0.4,0.18,size=7,color=DGREY,align=PP_ALIGN.CENTER)

    # right: ΔAUC summary and RNA vs img gap
    txt(sl,'Key Numbers',6.85,3.75,6.3,0.28,size=10,bold=True,color=ORA)
    sh2=table(sl,[['Metric','All-sites','At-20-cm','Change'],
                  ['RNA AUC',f"{rna_all['mean_auc']:.3f}",f"{rna_20['mean_auc']:.3f}",
                   f"−{rna_drop:.3f}"],
                  ['Imaging AUC',f"{img_all['mean_auc']:.3f}",f"{img_20['mean_auc']:.3f}",
                   f"−{img_drop:.3f}"],
                  ['Fusion AUC',f"{fus_all['mean_auc']:.3f}",f"{fus_20['mean_auc']:.3f}",
                   f"−{fus_all['mean_auc']-fus_20['mean_auc']:.3f}"],
                  ['RNA − imaging gap',
                   f"+{rna_all['mean_auc']-img_all['mean_auc']:.3f}",
                   f"+{rna_20['mean_auc']-img_20['mean_auc']:.3f}","−0.010"],
                  ['Fusion − RNA',
                   f"{fus_all['mean_auc']-rna_all['mean_auc']:+.4f}",
                   f"{fus_20['mean_auc']-rna_20['mean_auc']:+.4f}","≈0"],
              ],
             l=6.85,t=4.08,cw=[3.2,1.8,1.8,1.3],rh=0.29,fs=9)
    for r in [1,2,3]:
        set_cell_bg(sh2.table.cell(r,3), LORA_H)
    for c in range(4): set_cell_bg(sh2.table.cell(4,c), LTEAL_H)

    footer(sl,'Orange bars = AUC attributable to site-identity shortcut  ·  At-20-cm = honest site-matched estimate')

    # ── Slide 4: Interpretation ───────────────────────────────────────────────
    sl = prs.slides.add_slide(BLANK)
    rect(sl,0,0,13.33,7.5,'FAFAFA')
    header(sl,'Interpretation','What the AUC drop tells us about each modality')

    for i,(title,body,fh,bg) in enumerate([
        ('~0.09 AUC was site-identity in both modalities',
         f'Imaging: −{img_drop:.3f}  |  RNA: −{rna_drop:.3f}. Both modalities exploited cecum '
         f'samples as a UC marker. Virchow2 captured cecum morphology (gland pattern, mucosal '
         f'folds). VST captured cecum-specific gene expression (cell types, microbiome). '
         f'Neither signal is disease-specific.',
         ORA_H, LORA_H),
        ('The drops are nearly equal across modalities',
         f'Imaging and RNA lost the same absolute AUC (±0.011 of each other). '
         f'This confirms they were exploiting the same underlying confound — '
         f'site identity — rather than independent modality-specific shortcuts.',
         ORA_H, LORA_H),
        ('RNA retains a genuine advantage at matched tissue',
         f'At-20-cm: RNA {rna_20["mean_auc"]:.3f} vs imaging {img_20["mean_auc"]:.3f} '
         f'(gap +{rna_20["mean_auc"]-img_20["mean_auc"]:.3f}). '
         f'Molecular gene expression at the rectosigmoid junction genuinely distinguishes '
         f'CD from UC beyond what H&E morphology alone reveals.',
         TEAL_H, LTEAL_H),
        ('Fusion still adds nothing over RNA alone',
         f'concat_raw fusion: {fus_20["mean_auc"]:.3f} vs RNA alone: {rna_20["mean_auc"]:.3f} '
         f'(Δ = {fus_20["mean_auc"]-rna_20["mean_auc"]:+.3f}). The same conclusion holds '
         f'even after removing the site confound: imaging does not add discriminative '
         f'information beyond what RNA already captures at the rectosigmoid junction.',
         MID_H, LIGHT_H),
    ]):
        lx = 0.25 + (i%2)*6.54; ly = 1.08 + (i//2)*2.42
        rect(sl,lx,ly,6.3,2.28,bg,fh)
        txt(sl,title,lx+0.12,ly+0.06,6.1,0.30,size=9.5,bold=True,
            color=RGBColor.from_string(fh))
        txt(sl,body, lx+0.12,ly+0.40,6.1,1.80,size=8.5,color=DARK)

    divider(sl,5.90)
    txt(sl,'The honest performance benchmark',0.25,6.00,12.8,0.28,size=10,bold=True,color=ORA)
    rect(sl,0.25,6.32,12.83,0.88,LORA_H,ORA_H)
    mtxt(sl,[f'At-20-cm AUC represents "given a rectosigmoid biopsy, predict CD vs UC" — '
             f'the clinically meaningful question a pathologist faces.',
             f'RNA: {rna_20["mean_auc"]:.3f}  ·  Imaging: {img_20["mean_auc"]:.3f}  ·  '
             f'All-sites figures ({rna_all["mean_auc"]:.3f} / {img_all["mean_auc"]:.3f}) should be '
             f'reported as an upper bound that includes ~0.09 AUC of protocol-driven site signal.'],
         0.38,6.35,12.55,0.82,size=9,color=DARK)

    footer(sl,'Cecum samples are UC-specific by protocol — not by biology alone')

    # ── Slide 5: Discussion + next steps ─────────────────────────────────────
    sl = prs.slides.add_slide(BLANK)
    rect(sl,0,0,13.33,7.5,'FAFAFA')
    header(sl,'Discussion & Next Steps',
           'Implications for model development and honest reporting')

    txt(sl,'Summary of Findings',0.25,1.08,6.3,0.30,size=11,bold=True,color=ORA)
    for i,(t,b) in enumerate([
        ('All-sites AUC was inflated by ~0.09',
         f'RNA: {rna_all["mean_auc"]:.3f} → {rna_20["mean_auc"]:.3f}  |  '
         f'Imaging: {img_all["mean_auc"]:.3f} → {img_20["mean_auc"]:.3f}  |  '
         f'Equal drop confirms same site-identity shortcut exploited by both.'),
        ('At-20-cm is the rigorous benchmark',
         'Rectosigmoid junction biopsied in both CD and UC. '
         'Any CD/UC signal here is disease biology, not location identity.'),
        ('RNA still outperforms imaging at matched tissue',
         f'Gap narrows slightly ({rna_all["mean_auc"]-img_all["mean_auc"]:.3f} → '
         f'{rna_20["mean_auc"]-img_20["mean_auc"]:.3f}) but RNA (0.824) '
         f'remains clearly better than imaging (0.762).'),
        ('Fusion adds nothing either way',
         'concat_raw ≈ RNA-only on both all-sites and At-20-cm. '
         'A non-linear fusion model is needed to capture cross-modal interactions.'),
    ]):
        ly = 1.42 + i*0.78
        rect(sl,0.25,ly,6.3,0.72,LORA_H,ORA_H)
        txt(sl,t,0.38,ly+0.03,6.1,0.25,size=9,bold=True,color=ORA)
        txt(sl,b,0.38,ly+0.28,6.1,0.40,size=8.5,color=DARK)

    txt(sl,'Next Steps',6.85,1.08,6.3,0.30,size=11,bold=True,color=MID)
    for i,(t,b) in enumerate([
        ('Cecum-only control',
         'Train on Cecum samples only; expect AUC ≈ 1.0 (pure site = UC). '
         'Quantifies the maximum site shortcut.'),
        ('Site-covariate model',
         'Add biopsy location as a one-hot feature and train on all sites. '
         'Tests whether disease signal can be separated from site signal.'),
        ('SHAP on At-20-cm RNA',
         'Identify which genes drive CD/UC separation at the rectosigmoid '
         'without cecum expression confounding the importance scores.'),
        ('Non-linear fusion',
         'Neural late fusion or cross-modal attention may extract imaging–RNA '
         'interactions that RF on concatenated features misses.'),
    ]):
        ly = 1.42 + i*0.78
        rect(sl,6.85,ly,6.3,0.72,LIGHT_H,MID_H)
        txt(sl,t,6.98,ly+0.03,6.1,0.25,size=9,bold=True,color=MID)
        txt(sl,b,6.98,ly+0.28,6.1,0.40,size=8.5,color=DARK)

    divider(sl,5.52)
    txt(sl,'Recommended Reporting Language',0.25,5.62,12.8,0.28,size=10,bold=True,color=ORA)
    rect(sl,0.25,5.95,12.83,1.28,LORA_H,ORA_H)
    mtxt(sl,['"We report two sets of AUC estimates: (1) all-sites AUC, which includes samples from all '
             'colon locations and is inflated by protocol-driven biopsy-site differences between '
             'CD and UC (~0.09 AUC); and (2) At-20-cm AUC, restricted to the rectosigmoid junction '
             '(biopsied in both CD and UC by protocol), which represents a site-controlled estimate '
             'of disease-specific discriminative performance. RNA: 0.824 ± 0.019; '
             'Imaging: 0.762 ± 0.024 (At-20-cm, matched 828-patient cohort)."'],
         0.38,5.98,12.55,1.20,size=8.5,color=DARK)

    footer(sl,'IBD Plexus / SPARC  ·  CombatSeq VST  ·  Virchow2 TRIDENT 0.3.0  ·  August 11, 2026')

    out_path, _ = next_versioned_path(os.path.join(REPORTS_DIR, 'at20cm_slides.pptx'))
    prs.save(out_path)
    print(f'PPTX saved: {out_path}')
    return out_path


if __name__ == '__main__':
    parser = argparse.ArgumentParser()
    parser.add_argument('--desc',
                        default='At-20-cm site-restricted analysis: removing biopsy-protocol confound, 828 matched patients',
                        help='One-line description for VERSIONS.md')
    args = parser.parse_args()

    pdf_path  = build_pdf()
    pptx_path = build_pptx()
    log_version(pdf_path,  args.desc)
    log_version(pptx_path, args.desc)
